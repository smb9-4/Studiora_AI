"""
Professor AI — Teaching Assistant Backend (RAG + Auth + Per-Chat Isolation)
============================================================
Adapted from the Student AI backend. Same core architecture — per-chat
Chroma isolation, SQLite-persisted history, long-term memory with
importance-weighted recall, the structured-JSON document generation engine —
with the Professor-specific constraints from the module spec layered on top:

  - Long-term memory is restricted to academic-preference categories only
    (name, department, subjects taught, teaching style, question pattern,
    output format, frequently used prompts). No general personal facts, no
    conversation content, no student data are ever extracted into memory.
  - The system prompt enforces a professional, fact-based, education-only
    tone: no emotional/motivational counseling, no fabricated student
    records, no code execution, no going outside the teaching domain.
  - Document generation is extended with faculty content types: lecture
    notes, assignments, quizzes, MCQ sets, descriptive question sets, lab
    manuals, question papers, answer keys, lesson plans, teaching
    schedules, and rubric-based evaluation reports.
  - Document intelligence adds .pptx extraction (python-pptx) alongside the
    existing PDF/DOCX/TXT support, per the spec's "Analyze uploaded PDFs,
    DOCX, PPTs, and TXT files."
  - Storage (SQLite DB, Chroma dir, uploads dir, generated-documents dir) is
    kept in entirely separate files/folders from the Student AI app, so the
    two can run side by side on the same machine with zero chance of one
    professor's or one student's data ever being visible to the other app.

Requirements: pip install flask flask-cors ollama chromadb pymupdf python-docx python-pptx reportlab werkzeug

Ollama models needed:
    ollama pull qwen2.5:7b
    ollama pull nomic-embed-text
    ollama pull qwen2.5vl:7b      # optional — image analysis (e.g. scanned/handwritten answer sheets)

Folder structure required:
    app.py
    templates/
        index.html
        login.html
    static/
        style.css
        script.js

Run: python app.py
Then open: http://localhost:5001   (different default port from the Student
                                     app so both can run at once — change
                                     PORT below if you'd rather pick your own)

ISOLATION MODEL (same as Student AI, unit = the chat)
------------------------------------------------------
  - Each chat has its own Chroma collection ("chat_<id>") for uploaded
    course material — a document uploaded in one chat is never visible to
    another chat, and never visible to another professor's account.
  - Each chat has its own message history and uploaded-file list in SQLite.
  - Long-term memory (academic preferences only) is scoped to the
    PROFESSOR's account and stored in its own per-professor Chroma
    collection, entirely separate from any chat's document collection.

FRONTEND CONTRACT (identical route shapes to the Student AI, so the same
frontend patterns/script.js structure carry over)
-------------------------------------------------------------------------
        GET    /api/chats
        POST   /api/chats
        PATCH  /api/chats/<chat_id>
        DELETE /api/chats/<chat_id>
        POST   /api/chats/<chat_id>/upload
        GET    /api/chats/<chat_id>/files
        DELETE /api/chats/<chat_id>/files/<filename>
        POST   /api/chats/<chat_id>/files/<filename>/reindex
        GET    /api/chats/<chat_id>/history
        DELETE /api/chats/<chat_id>/history
        POST   /api/chats/<chat_id>/chat        (streaming SSE)
        POST   /api/chats/<chat_id>/chat-image   (streaming SSE — optional vision model)
        GET    /api/memory
        POST   /api/memory
        DELETE /api/memory/<id>
        DELETE /api/memory
        POST   /api/generate-document
        GET    /api/documents
        GET    /download/<filename>
"""

import os
import re
import json
import time
import base64
import threading
import sqlite3
from pathlib import Path
from functools import wraps
from collections import Counter

import fitz  # PyMuPDF
import ollama
import chromadb
from chromadb.utils import embedding_functions
from flask import (Flask, request, jsonify, session, Response, render_template,
                    redirect, url_for, send_from_directory)
from flask_cors import CORS
from werkzeug.security import generate_password_hash, check_password_hash

# ── optional docx support (reading uploads AND writing generated documents) ──
try:
    from docx import Document as DocxDocument
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

# ── optional pptx support (reading uploaded lecture slides) — Professor AI
# spec requires PPT analysis; the Student AI backend never needed this. ──
try:
    from pptx import Presentation
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

# ── optional PDF-generation support (for the document generation engine) ──
try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                     TableStyle, ListFlowable, ListItem)
    PDF_OK = True
except ImportError:
    PDF_OK = False

# ═══════════════════════════════════════════════════════
# CONFIG
# ═══════════════════════════════════════════════════════
CHAT_MODEL = "qwen2.5:7b"
EMBED_MODEL = "nomic-embed-text"
CHUNK_SIZE = 512
CHUNK_OVERLAP = 64
TOP_K = 5
MAX_HISTORY = 20  # number of (user+assistant) turn-pairs kept for prompt context
EMBED_BATCH_SIZE = 32  # chunks per embedding call — large single batches crash the Ollama runner

# ── long-term memory (cross-chat) ──
MEMORY_MODEL = CHAT_MODEL          # model used to extract facts; can point at a smaller/faster model
# New messages (user+assistant) accumulated in a chat before a background
# extraction run fires. Was 6 (i.e. waited for 3 full turns) — that batching
# saved LLM calls, but it also meant a fact stated in this chat could be
# invisible to every OTHER chat for a while: if you stated your interests
# and then immediately opened a new chat and asked about them, extraction
# for the first chat often hadn't run yet, so the fact was never even in
# storage yet for the new chat to retrieve. 2 means extraction fires after
# every single exchange — it's still a background thread (doesn't block the
# reply you're reading), so the only cost is more frequent small LLM calls,
# which is the right trade for a local single-user assistant where getting
# a stated fact into memory promptly matters more than saving a few calls.
MEMORY_EXTRACTION_THRESHOLD = 2
MEMORY_TOP_K = 4                   # long-term facts injected into the prompt per turn
MEMORY_CANDIDATE_POOL = 15         # candidates pulled from Chroma before importance-weighted re-rank
MEMORY_EXCERPT_CHAR_CAP = 8000     # cap on how much raw conversation text goes into one extraction call
# Professor AI memory is DELIBERATELY narrower than the Student AI's. Per the
# spec's "Professor Memory" section, only academic-preference facts may be
# retained long-term — never personal conversations, sensitive information,
# student private data, or general chit-chat. Each category maps directly to
# one line of the spec's "Store only" list:
#   identity          -> Name
#   department         -> Department
#   subject            -> Subjects taught
#   teaching_style      -> Preferred teaching style
#   question_pattern    -> Preferred question pattern
#   output_format       -> Preferred output format
#   frequent_prompt     -> Frequently used prompts
# ("Recently uploaded documents" and "Recent AI sessions" are already tracked
# structurally via chat_files / chats and don't need a memory category.)
MEMORY_CATEGORIES = ("identity", "department", "subject", "teaching_style",
                      "question_pattern", "output_format", "frequent_prompt")
MEMORY_MIN_IMPORTANCE = 0.35       # extracted facts scored below this are discarded as not worth remembering

# Two-band distance model for deciding what happens when a new fact is close
# to something already stored (nearest-neighbour Chroma distance, smaller = closer):
#   dist < MEMORY_DEDUP_DISTANCE   -> true duplicate, skip (or reinforce importance)
#   dist < MEMORY_UPDATE_DISTANCE  -> same topic, info changed -> UPDATE the existing
#                                      row in place instead of stacking a near-duplicate
#                                      (this is what lets "uses React" become "prefers
#                                      vanilla JS" instead of just sitting there twice)
#   dist >= MEMORY_UPDATE_DISTANCE -> genuinely new fact -> insert
MEMORY_DEDUP_DISTANCE = 0.10
MEMORY_UPDATE_DISTANCE = 0.32
MEMORY_RETRIEVAL_IMPORTANCE_WEIGHT = 0.25  # blend weight for importance vs semantic closeness on recall

# How close (Chroma distance, smaller = closer) a stored fact must be to the
# CURRENT message before it's allowed into the prompt at all. Without this,
# retrieve_memory() always returns its top-k regardless of relevance, which is
# what caused "hii" / "hello good evening" to drag in "uses RTX 4050, 6GB VRAM"
# every single time — the memory collection is tiny, so *something* is always
# the nearest neighbour even when nothing is actually related.
#
# retrieve_memory() and store_memory_fact() now embed through _memory_embed()
# above, which applies nomic-embed-text's "search_query: " / "search_document: "
# task prefixes — the asymmetric pairing the model is actually trained for.
# That gives much cleaner separation than comparing both sides unprefixed, so
# 0.55 is a reasonable starting point here (tighter than the old 0.85, which
# was set back when both sides were unprefixed and everything looked "close").
# retrieve_memory() prints the real distance on every call — watch the
# console and tune from what you actually see if greetings still leak facts.
MEMORY_MAX_DISTANCE = 0.55

# ── document RAG retrieval quality ──
RAG_CANDIDATE_MULTIPLIER = 4   # pull TOP_K * this many candidates from Chroma before re-ranking
RAG_DISTANCE_MARGIN = 0.40     # keep candidates within this margin of the single best distance...
RAG_MAX_DISTANCE = 1.60        # ...or under this absolute ceiling, whichever is more permissive
RAG_LEXICAL_WEIGHT = 0.30      # blend weight for keyword overlap vs semantic closeness (0 = pure embeddings)

# ── document generation engine (Phase 1: prompt -> structured content -> DOCX/PDF) ──
# The standalone /api/generate-document endpoint stays prompt-only per the
# original Phase 1 spec (no chat_id, nothing to ground it in). The
# chat-triggered path (handle_chat_document_generation) is different: it DOES
# have a chat_id, and that chat may have an uploaded PDF sitting right there
# in Chroma, so "summarize the important points from the PDF" now pulls the
# actual PDF content in as source material instead of asking the LLM to
# invent a plausible-looking summary from six words of request text.
DOC_GEN_MODEL = CHAT_MODEL
DOC_MAX_SECTIONS = 20           # sanity cap on how many sections one document can have
DOC_ALLOWED_FORMATS = ("pdf", "docx", "both")
DOC_GEN_RAG_K = 15              # wider net than normal chat QA (TOP_K=5) — summarizing
                                 # a whole document needs coverage across it, not just
                                 # the handful of chunks nearest one narrow query
# Ollama's chat() defaults to a fairly small num_predict if it isn't set
# explicitly, which was silently truncating generated documents down to
# 1-2 pages regardless of how many sections were requested. This overrides
# that ceiling for the content-generation call specifically — it's generous
# on purpose since a cut-off document is worse than a slightly slow one.
DOC_GEN_MAX_TOKENS = 7000
GENERATED_DOCS_DIR = Path("generated_documents_professor")
GENERATED_DOCS_DIR.mkdir(exist_ok=True)

# ── vision (optional — e.g. scanned/handwritten answer sheets, exam papers) ──
# Point this at whatever tag `ollama list` shows for your Qwen2.5-VL pull.
# If you don't have a vision model pulled, just leave the Attach-image button
# unused — every other feature works fine without it.
VISION_MODEL = "qwen2.5vl:7b"
VISION_ALLOWED_EXT = (".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp")
VISION_MAX_BYTES = 15 * 1024 * 1024
VISION_MAX_TOKENS = 2000

# Entirely separate storage from the Student AI app — different DB file,
# different Chroma directory, different uploads/generated-docs folders. This
# is what actually guarantees "Documents are isolated per account. AI
# responses must never expose another user's data": even if both apps run
# from the same working directory on the same machine, there is no shared
# file, table, or collection either one could accidentally read from.
UPLOAD_DIR = Path("uploads_professor")
CHROMA_DIR = Path("chroma_db_professor")
DB_PATH = Path("professors.db")

UPLOAD_DIR.mkdir(exist_ok=True)
CHROMA_DIR.mkdir(exist_ok=True)

# ═══════════════════════════════════════════════════════
# FLASK SETUP
# ═══════════════════════════════════════════════════════
app = Flask(__name__, static_folder="static", static_url_path="/static")
app.secret_key = os.environ.get("PROFESSOR_AI_SECRET", "dev-secret-change-in-production-" + os.urandom(8).hex())
app.config["SESSION_COOKIE_SAMESITE"] = "Lax"
CORS(app, supports_credentials=True)

# ═══════════════════════════════════════════════════════
# DATABASE (SQLite) — users, chats, messages, files
# ═══════════════════════════════════════════════════════
def get_db():
    conn = sqlite3.connect(str(DB_PATH))
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA foreign_keys = ON")
    return conn


def init_db():
    conn = get_db()
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            name TEXT NOT NULL,
            password_hash TEXT NOT NULL,
            role TEXT DEFAULT 'professor',
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP
        );

        CREATE TABLE IF NOT EXISTS chats (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL REFERENCES users(id) ON DELETE CASCADE,
            title TEXT NOT NULL DEFAULT 'New chat',
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            updated_at DATETIME DEFAULT CURRENT_TIMESTAMP
        );

        CREATE TABLE IF NOT EXISTS messages (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            chat_id INTEGER NOT NULL REFERENCES chats(id) ON DELETE CASCADE,
            role TEXT NOT NULL,
            content TEXT NOT NULL,
            image_path TEXT,
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP
        );

        CREATE TABLE IF NOT EXISTS chat_files (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            chat_id INTEGER NOT NULL REFERENCES chats(id) ON DELETE CASCADE,
            filename TEXT NOT NULL,
            stored_name TEXT NOT NULL,
            chunks INTEGER NOT NULL,
            pages INTEGER NOT NULL,
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            UNIQUE(chat_id, filename)
        );

        -- Long-term, cross-chat memory: atomic facts extracted from conversations.
        -- Scoped to the USER (not the chat) on purpose — this is the one place
        -- information is intentionally allowed to cross chat boundaries.
        CREATE TABLE IF NOT EXISTS user_memory (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL REFERENCES users(id) ON DELETE CASCADE,
            memory_text TEXT NOT NULL,
            category TEXT NOT NULL DEFAULT 'auto',   -- identity | preference | project | skill | goal | manual
            importance REAL NOT NULL DEFAULT 0.7,    -- 0..1, how much this should weigh during recall
            confidence REAL NOT NULL DEFAULT 0.8,    -- 0..1, how sure we are this fact is accurate/current
            source_chat_id INTEGER REFERENCES chats(id) ON DELETE SET NULL,
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            updated_at DATETIME DEFAULT CURRENT_TIMESTAMP
        );

        -- Watermark of how far each chat's messages have been scanned for
        -- memory extraction, so re-runs don't reprocess the same turns.
        CREATE TABLE IF NOT EXISTS memory_extraction_log (
            chat_id INTEGER PRIMARY KEY REFERENCES chats(id) ON DELETE CASCADE,
            last_extracted_message_id INTEGER NOT NULL DEFAULT 0,
            updated_at DATETIME DEFAULT CURRENT_TIMESTAMP
        );

        -- Record of documents generated by the Phase 1 document engine, so a
        -- user can see/re-download what they've generated after leaving the page.
        CREATE TABLE IF NOT EXISTS generated_documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL REFERENCES users(id) ON DELETE CASCADE,
            title TEXT NOT NULL,
            doc_type TEXT NOT NULL DEFAULT 'notes',
            filename TEXT NOT NULL,
            format TEXT NOT NULL,
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP
        );

        CREATE INDEX IF NOT EXISTS idx_chats_user ON chats(user_id);
        CREATE INDEX IF NOT EXISTS idx_messages_chat ON messages(chat_id);
        CREATE INDEX IF NOT EXISTS idx_files_chat ON chat_files(chat_id);
        CREATE INDEX IF NOT EXISTS idx_memory_user ON user_memory(user_id);
        CREATE INDEX IF NOT EXISTS idx_gendocs_user ON generated_documents(user_id);
    """)

    # Migration for DBs created before importance/confidence/updated_at existed —
    # CREATE TABLE IF NOT EXISTS above won't retroactively add columns to a table
    # that's already there, so patch it in by hand if missing.
    existing_cols = {row["name"] for row in conn.execute("PRAGMA table_info(user_memory)")}
    if "importance" not in existing_cols:
        conn.execute("ALTER TABLE user_memory ADD COLUMN importance REAL NOT NULL DEFAULT 0.7")
    if "confidence" not in existing_cols:
        conn.execute("ALTER TABLE user_memory ADD COLUMN confidence REAL NOT NULL DEFAULT 0.8")
    if "updated_at" not in existing_cols:
        # SQLite refuses ADD COLUMN with a non-constant default like
        # CURRENT_TIMESTAMP (that's only allowed at CREATE TABLE time), so add
        # it bare and backfill existing rows from created_at instead.
        conn.execute("ALTER TABLE user_memory ADD COLUMN updated_at DATETIME")
        conn.execute("UPDATE user_memory SET updated_at = created_at WHERE updated_at IS NULL")

    msg_cols = {row["name"] for row in conn.execute("PRAGMA table_info(messages)")}
    if "image_path" not in msg_cols:
        conn.execute("ALTER TABLE messages ADD COLUMN image_path TEXT")

    conn.commit()
    conn.close()


init_db()


# ── user helpers ──
def get_user(username: str):
    conn = get_db()
    row = conn.execute("SELECT * FROM users WHERE username = ?", (username,)).fetchone()
    conn.close()
    return dict(row) if row else None


def create_user(username: str, name: str, password: str, role: str = "professor"):
    conn = get_db()
    try:
        conn.execute(
            "INSERT INTO users (username, name, password_hash, role) VALUES (?, ?, ?, ?)",
            (username, name, generate_password_hash(password), role)
        )
        conn.commit()
        return True
    except sqlite3.IntegrityError:
        return False
    finally:
        conn.close()


# ── chat helpers ──
def create_chat(user_id: int, title: str = "New chat") -> int:
    conn = get_db()
    cur = conn.execute("INSERT INTO chats (user_id, title) VALUES (?, ?)", (user_id, title))
    conn.commit()
    chat_id = cur.lastrowid
    conn.close()
    return chat_id


def list_chats(user_id: int):
    conn = get_db()
    rows = conn.execute(
        "SELECT id, title, created_at, updated_at FROM chats WHERE user_id = ? ORDER BY updated_at DESC",
        (user_id,)
    ).fetchall()
    conn.close()
    return [dict(r) for r in rows]


def get_chat(chat_id: int):
    conn = get_db()
    row = conn.execute("SELECT * FROM chats WHERE id = ?", (chat_id,)).fetchone()
    conn.close()
    return dict(row) if row else None


def touch_chat(chat_id: int):
    conn = get_db()
    conn.execute("UPDATE chats SET updated_at = CURRENT_TIMESTAMP WHERE id = ?", (chat_id,))
    conn.commit()
    conn.close()


def rename_chat(chat_id: int, title: str):
    conn = get_db()
    conn.execute("UPDATE chats SET title = ? WHERE id = ?", (title, chat_id))
    conn.commit()
    conn.close()


def delete_chat_row(chat_id: int):
    conn = get_db()
    conn.execute("DELETE FROM chats WHERE id = ?", (chat_id,))  # cascades to messages, chat_files
    conn.commit()
    conn.close()


def owns_chat(chat_id: int, username: str) -> bool:
    chat = get_chat(chat_id)
    if not chat:
        return False
    user = get_user(username)
    return bool(user) and chat["user_id"] == user["id"]


# ── message helpers ──
def add_message(chat_id: int, role: str, content: str, image_path: str | None = None):
    conn = get_db()
    conn.execute(
        "INSERT INTO messages (chat_id, role, content, image_path) VALUES (?, ?, ?, ?)",
        (chat_id, role, content, image_path)
    )
    conn.commit()
    conn.close()


def get_history(chat_id: int, limit: int = MAX_HISTORY):
    """Plain role+content history for feeding back into the LLM as prompt
    context. Deliberately excludes image_path — past images are never
    replayed into later turns. See get_history_for_display for the
    UI-facing version that includes them."""
    conn = get_db()
    rows = conn.execute(
        "SELECT role, content FROM messages WHERE chat_id = ? ORDER BY id DESC LIMIT ?",
        (chat_id, limit * 2)
    ).fetchall()
    conn.close()
    return [{"role": r["role"], "content": r["content"]} for r in reversed(rows)]


def get_history_for_display(chat_id: int, limit: int = MAX_HISTORY):
    """Same as get_history but includes image_path, so an attached image
    still shows as a thumbnail after a page reload."""
    conn = get_db()
    rows = conn.execute(
        "SELECT role, content, image_path FROM messages WHERE chat_id = ? ORDER BY id DESC LIMIT ?",
        (chat_id, limit * 2)
    ).fetchall()
    conn.close()
    return [{"role": r["role"], "content": r["content"], "image_path": r["image_path"]} for r in reversed(rows)]


def clear_chat_history(chat_id: int):
    conn = get_db()
    conn.execute("DELETE FROM messages WHERE chat_id = ?", (chat_id,))
    conn.commit()
    conn.close()


# ── file helpers ──
def record_file(chat_id: int, filename: str, stored_name: str, chunks: int, pages: int):
    conn = get_db()
    conn.execute("""
        INSERT INTO chat_files (chat_id, filename, stored_name, chunks, pages)
        VALUES (?, ?, ?, ?, ?)
        ON CONFLICT(chat_id, filename) DO UPDATE SET
            stored_name=excluded.stored_name, chunks=excluded.chunks, pages=excluded.pages
    """, (chat_id, filename, stored_name, chunks, pages))
    conn.commit()
    conn.close()


def list_chat_files(chat_id: int):
    conn = get_db()
    rows = conn.execute(
        "SELECT filename, stored_name, chunks, pages FROM chat_files WHERE chat_id = ?", (chat_id,)
    ).fetchall()
    conn.close()
    return [dict(r) for r in rows]


def get_chat_file(chat_id: int, filename: str):
    conn = get_db()
    row = conn.execute(
        "SELECT * FROM chat_files WHERE chat_id = ? AND filename = ?", (chat_id, filename)
    ).fetchone()
    conn.close()
    return dict(row) if row else None


def remove_chat_file(chat_id: int, filename: str):
    conn = get_db()
    conn.execute("DELETE FROM chat_files WHERE chat_id = ? AND filename = ?", (chat_id, filename))
    conn.commit()
    conn.close()


# ── long-term memory helpers (SQLite side; Chroma side is in the next section) ──
def list_user_memory(user_id: int):
    conn = get_db()
    rows = conn.execute(
        "SELECT id, memory_text, category, importance, confidence, source_chat_id, created_at, updated_at "
        "FROM user_memory WHERE user_id = ? ORDER BY importance DESC, created_at DESC",
        (user_id,)
    ).fetchall()
    conn.close()
    return [dict(r) for r in rows]


def get_memory_row(user_id: int, memory_id: int):
    conn = get_db()
    row = conn.execute(
        "SELECT id, memory_text, category, importance, confidence FROM user_memory WHERE id = ? AND user_id = ?",
        (memory_id, user_id)
    ).fetchone()
    conn.close()
    return dict(row) if row else None


def insert_memory_row(user_id: int, text: str, category: str, source_chat_id,
                       importance: float = 0.7, confidence: float = 0.8):
    conn = get_db()
    cur = conn.execute(
        "INSERT INTO user_memory (user_id, memory_text, category, importance, confidence, source_chat_id) "
        "VALUES (?, ?, ?, ?, ?, ?)",
        (user_id, text, category, importance, confidence, source_chat_id)
    )
    conn.commit()
    mem_id = cur.lastrowid
    conn.close()
    return mem_id


def update_memory_row(user_id: int, memory_id: int, text: str, category: str,
                       importance: float, confidence: float) -> bool:
    """Overwrites an existing memory's content in place — used when a new fact
    turns out to be an evolved version of something already stored (e.g. the
    user's tech stack changed), rather than a genuinely new, separate fact."""
    conn = get_db()
    cur = conn.execute(
        "UPDATE user_memory SET memory_text = ?, category = ?, importance = ?, confidence = ?, "
        "updated_at = CURRENT_TIMESTAMP WHERE id = ? AND user_id = ?",
        (text, category, importance, confidence, memory_id, user_id)
    )
    conn.commit()
    ok = cur.rowcount > 0
    conn.close()
    return ok


def delete_memory_row(user_id: int, memory_id: int) -> bool:
    conn = get_db()
    row = conn.execute(
        "SELECT id FROM user_memory WHERE id = ? AND user_id = ?", (memory_id, user_id)
    ).fetchone()
    if not row:
        conn.close()
        return False
    conn.execute("DELETE FROM user_memory WHERE id = ?", (memory_id,))
    conn.commit()
    conn.close()
    return True


def clear_user_memory_rows(user_id: int):
    conn = get_db()
    conn.execute("DELETE FROM user_memory WHERE user_id = ?", (user_id,))
    conn.commit()
    conn.close()


def get_extraction_watermark(chat_id: int) -> int:
    conn = get_db()
    row = conn.execute(
        "SELECT last_extracted_message_id FROM memory_extraction_log WHERE chat_id = ?", (chat_id,)
    ).fetchone()
    conn.close()
    return row["last_extracted_message_id"] if row else 0


def set_extraction_watermark(chat_id: int, last_message_id: int):
    conn = get_db()
    conn.execute("""
        INSERT INTO memory_extraction_log (chat_id, last_extracted_message_id, updated_at)
        VALUES (?, ?, CURRENT_TIMESTAMP)
        ON CONFLICT(chat_id) DO UPDATE SET
            last_extracted_message_id = excluded.last_extracted_message_id,
            updated_at = CURRENT_TIMESTAMP
    """, (chat_id, last_message_id))
    conn.commit()
    conn.close()


def count_new_messages_since_watermark(chat_id: int) -> int:
    last_id = get_extraction_watermark(chat_id)
    conn = get_db()
    row = conn.execute(
        "SELECT COUNT(*) AS c FROM messages WHERE chat_id = ? AND id > ?", (chat_id, last_id)
    ).fetchone()
    conn.close()
    return row["c"]


def get_new_messages_since_watermark(chat_id: int):
    last_id = get_extraction_watermark(chat_id)
    conn = get_db()
    rows = conn.execute(
        "SELECT id, role, content FROM messages WHERE chat_id = ? AND id > ? ORDER BY id ASC",
        (chat_id, last_id)
    ).fetchall()
    conn.close()
    return [dict(r) for r in rows]


# ═══════════════════════════════════════════════════════
# AUTH DECORATOR
# ═══════════════════════════════════════════════════════
def login_required(f):
    @wraps(f)
    def wrapper(*args, **kwargs):
        if "username" not in session:
            if request.path.startswith("/api/"):
                return jsonify({"error": "Not authenticated"}), 401
            return redirect(url_for("login_page"))
        return f(*args, **kwargs)
    return wrapper


def chat_access_required(f):
    """Ensures the chat_id in the URL belongs to the logged-in user."""
    @wraps(f)
    def wrapper(chat_id, *args, **kwargs):
        chat_id = int(chat_id)
        if not owns_chat(chat_id, session["username"]):
            return jsonify({"error": "Chat not found"}), 404
        return f(chat_id, *args, **kwargs)
    return wrapper


# ═══════════════════════════════════════════════════════
# CHROMADB + EMBEDDINGS — one collection per CHAT, not per user
# ═══════════════════════════════════════════════════════
chroma_client = chromadb.PersistentClient(path=str(CHROMA_DIR))
ollama_ef = embedding_functions.OllamaEmbeddingFunction(
    url="http://localhost:11434/api/embeddings",
    model_name=EMBED_MODEL,
)


def _memory_embed(text: str, is_query: bool) -> list[float]:
    """Embed text for the long-term-memory collection using nomic-embed-text's
    documented task prefixes ('search_query: ' for the live query side,
    'search_document: ' for the stored-fact side). nomic-embed-text is
    trained expecting these prefixes for asymmetric retrieval — without
    them, short texts (a two-word greeting vs. a one-line stored fact) can
    sit closer together in embedding space than they semantically should,
    which is the real reason irrelevant memories kept surfacing on plain
    greetings no matter where the distance cutoff was set. Scoped to memory
    only — chat_<id> document collections keep using the plain, unprefixed
    embedding_function below, so document RAG is completely unaffected."""
    prefix = "search_query: " if is_query else "search_document: "
    resp = ollama.embeddings(model=EMBED_MODEL, prompt=prefix + text)
    return resp["embedding"]

_COLLECTION_CACHE: dict = {}
_COLLECTION_LOCK = threading.Lock()


def get_chat_collection(chat_id: int):
    """Each chat gets its own isolated vector collection: chat_<id>."""
    col_name = f"chat_{chat_id}"
    with _COLLECTION_LOCK:
        if col_name not in _COLLECTION_CACHE:
            try:
                col = chroma_client.get_collection(col_name, embedding_function=ollama_ef)
            except Exception:
                col = chroma_client.create_collection(col_name, embedding_function=ollama_ef)
            _COLLECTION_CACHE[col_name] = col
        return _COLLECTION_CACHE[col_name]


def drop_chat_collection(chat_id: int):
    col_name = f"chat_{chat_id}"
    with _COLLECTION_LOCK:
        _COLLECTION_CACHE.pop(col_name, None)
    try:
        chroma_client.delete_collection(col_name)
    except Exception:
        pass


def get_user_memory_collection(user_id: int):
    """Each user gets one isolated long-term-memory collection: user_memory_<id>.
    Deliberately separate from chat_<id> collections (which hold uploaded-document
    chunks) — memory facts and document chunks are never retrieved together."""
    col_name = f"user_memory_{user_id}"
    with _COLLECTION_LOCK:
        if col_name not in _COLLECTION_CACHE:
            try:
                col = chroma_client.get_collection(col_name, embedding_function=ollama_ef)
            except Exception:
                col = chroma_client.create_collection(col_name, embedding_function=ollama_ef)
            _COLLECTION_CACHE[col_name] = col
        return _COLLECTION_CACHE[col_name]


def drop_user_memory_collection(user_id: int):
    col_name = f"user_memory_{user_id}"
    with _COLLECTION_LOCK:
        _COLLECTION_CACHE.pop(col_name, None)
    try:
        chroma_client.delete_collection(col_name)
    except Exception:
        pass


# ═══════════════════════════════════════════════════════
# TEXT EXTRACTION
# ═══════════════════════════════════════════════════════
def extract_pdf(path: Path) -> tuple[str, dict]:
    doc = fitz.open(str(path))
    pages, parts = {}, []
    for i, page in enumerate(doc, 1):
        t = page.get_text()
        pages[i] = t
        parts.append(t)
    return "\n".join(parts), pages


def extract_txt(path: Path) -> tuple[str, dict]:
    text = path.read_text(errors="ignore")
    lines = text.splitlines()
    page_size = 50
    pages = {}
    for i in range(0, len(lines), page_size):
        pg = i // page_size + 1
        pages[pg] = "\n".join(lines[i:i + page_size])
    return text, pages


def extract_docx(path: Path) -> tuple[str, dict]:
    if not DOCX_OK:
        raise RuntimeError("python-docx not installed")
    doc = DocxDocument(str(path))
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    page_size = 30
    pages = {}
    for i in range(0, len(paras), page_size):
        pg = i // page_size + 1
        pages[pg] = "\n".join(paras[i:i + page_size])
    return "\n".join(paras), pages


def extract_pptx(path: Path) -> tuple[str, dict]:
    """Extracts lecture-slide text for RAG. One "page" per slide (matching
    what a professor actually means by "slide 12") — pulls every text frame
    on the slide plus, where present, the speaker notes, since notes often
    carry the actual explanatory content a slide's bullet points only summarize."""
    if not PPTX_OK:
        raise RuntimeError("python-pptx not installed")
    prs = Presentation(str(path))
    pages, parts = {}, []
    for i, slide in enumerate(prs.slides, 1):
        chunks = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                chunks.append(shape.text_frame.text.strip())
            elif shape.has_table:
                for row in shape.table.rows:
                    chunks.append(" | ".join(cell.text.strip() for cell in row.cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            chunks.append(f"[Speaker notes] {slide.notes_slide.notes_text_frame.text.strip()}")
        slide_text = f"[Slide {i}]\n" + "\n".join(chunks) if chunks else f"[Slide {i}]"
        pages[i] = slide_text
        parts.append(slide_text)
    return "\n\n".join(parts), pages


def extract_file(path: Path) -> tuple[str, dict]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return extract_pdf(path)
    elif ext == ".txt":
        return extract_txt(path)
    elif ext == ".docx":
        return extract_docx(path)
    elif ext in (".pptx", ".ppt"):
        if ext == ".ppt":
            raise ValueError("Legacy .ppt isn't supported — please save as .pptx and re-upload.")
        return extract_pptx(path)
    else:
        raise ValueError(f"Unsupported format: {ext}")


# ═══════════════════════════════════════════════════════
# CHUNKING
# ═══════════════════════════════════════════════════════
def chunk_text(text: str, page_map: dict, filename: str, chat_id: int) -> list[dict]:
    chunks = []
    start, idx = 0, 0
    while start < len(text):
        end = min(start + CHUNK_SIZE, len(text))
        chunk = text[start:end]
        ratio = start / max(len(text), 1)
        pages = sorted(page_map.keys())
        page = pages[int(ratio * len(pages))] if pages else 1
        chunks.append({
            "text": chunk,
            # chunk_id namespaced by chat_id so the same filename in two
            # different chats never collides inside the shared Chroma store
            "chunk_id": f"chat{chat_id}__{filename}__chunk_{idx}",
            "filename": filename,
            "page": page,
        })
        start += CHUNK_SIZE - CHUNK_OVERLAP
        idx += 1
    return chunks


def add_chunks_in_batches(collection, ids: list, texts: list, metadatas: list,
                           batch_size: int = EMBED_BATCH_SIZE):
    """
    Send chunks to Chroma/Ollama in small batches instead of one giant call.
    A single collection.add() with thousands of texts (e.g. a 4MB+ PDF chunked
    at 512 chars) can send an oversized request to Ollama's embedding runner,
    which then drops the connection and every remaining chunk fails with it.
    Batching keeps each embedding request small and lets one bad batch fail
    without losing chunks that already succeeded.
    """
    total = len(ids)
    for start in range(0, total, batch_size):
        end = start + batch_size
        batch_ids       = ids[start:end]
        batch_texts     = texts[start:end]
        batch_metadatas = metadatas[start:end]
        last_err = None
        for attempt in range(3):
            try:
                collection.add(ids=batch_ids, documents=batch_texts, metadatas=batch_metadatas)
                last_err = None
                break
            except Exception as e:
                last_err = e
                time.sleep(1.5 * (attempt + 1))  # give the Ollama runner a moment to recover
        if last_err is not None:
            raise RuntimeError(
                f"Embedding failed on chunk batch {start}-{end} of {total} "
                f"(is Ollama running and reachable?): {last_err}"
            )


# ═══════════════════════════════════════════════════════
# RAG RETRIEVAL
# ═══════════════════════════════════════════════════════
# Hybrid retrieval: pure embedding similarity (what this had before) tends to
# miss exact terms — a formula, a proper noun, a specific section heading —
# that a plain keyword match would catch, and it always returns exactly k
# results even when nothing in the collection is actually relevant, which lets
# irrelevant chunks quietly leak into the model's context. This version pulls
# a larger semantic candidate pool, blends embedding closeness with lexical
# (keyword) overlap, drops candidates that aren't actually close to the query,
# and de-duplicates near-identical passages before returning the top k.

_STOPWORDS = set(
    "a an the is are was were be been being to of in on at for and or but with as by from "
    "this that these those it its into over under than then so such can could should would "
    "may might will shall do does did not no nor if while about above below between out up down "
    "i you he she we they what which who whom how why when where".split()
)


def _tokenize(text: str) -> set[str]:
    return {w for w in re.findall(r"[a-zA-Z0-9]+", (text or "").lower()) if len(w) > 1 and w not in _STOPWORDS}


def retrieve(collection, query: str, k: int = TOP_K) -> list[dict]:
    try:
        count = collection.count()
        if count == 0:
            return []

        pool_size = min(max(k * RAG_CANDIDATE_MULTIPLIER, 12), count)
        results = collection.query(
            query_texts=[query], n_results=pool_size,
            include=["documents", "metadatas", "distances"]
        )
        docs = results["documents"][0]
        metas = results["metadatas"][0]
        dists = results["distances"][0]
        if not docs:
            return []

        # Relative + absolute cutoff: keep anything close to the single best
        # match, OR anything under a generous absolute ceiling — whichever
        # keeps more of the genuinely relevant chunks. Only clearly weak
        # matches (far from the best hit AND far in absolute terms) get cut.
        best_dist = min(dists)
        relative_cutoff = best_dist + RAG_DISTANCE_MARGIN

        query_terms = _tokenize(query)
        scored = []
        for doc, meta, dist in zip(docs, metas, dists):
            if dist > relative_cutoff and dist > RAG_MAX_DISTANCE:
                continue
            closeness = 1.0 / (1.0 + dist)  # 0..1, higher = closer
            overlap = 0.0
            if query_terms:
                doc_terms = _tokenize(doc)
                if doc_terms:
                    overlap = len(query_terms & doc_terms) / len(query_terms)
            score = (1 - RAG_LEXICAL_WEIGHT) * closeness + RAG_LEXICAL_WEIGHT * overlap
            scored.append((score, doc, meta))

        scored.sort(key=lambda x: x[0], reverse=True)

        hits, seen = [], set()
        for score, doc, meta in scored:
            fingerprint = doc[:120].strip()
            if fingerprint in seen:
                continue  # skip near-identical/duplicate passages (e.g. overlapping chunks)
            seen.add(fingerprint)
            hits.append({"text": doc, "filename": meta.get("filename", "?"), "page": meta.get("page", "?")})
            if len(hits) >= k:
                break
        return hits
    except Exception:
        return []


# ═══════════════════════════════════════════════════════
# LONG-TERM MEMORY — storage, dedup, retrieval, extraction
# ═══════════════════════════════════════════════════════
#
# Design notes:
#   - Facts are stored ATOMICALLY (one short standalone sentence per row/vector),
#     not as a single rolling summary blob. This sidesteps the "merge two summaries
#     into one" problem entirely — new facts are just added, old facts age out
#     only if the user deletes them, and nothing ever gets silently rewritten or
#     lost to summarization drift.
#   - Storage is scoped to the USER, not the chat, since the whole point is that
#     it survives across chats. It lives in its own Chroma collection
#     (user_memory_<id>) and its own SQLite table, kept separate from per-chat
#     document collections and per-chat message history.
#   - A two-band nearest-neighbour check runs before every insert so a fact
#     isn't just blindly appended: a near-identical fact is treated as a
#     duplicate (skipped, or reinforced), while a fact that's clearly about
#     the same thing but with different content — same category, moderate
#     distance — UPDATES the existing row instead of creating a stale
#     duplicate. That's what lets memory actually evolve (e.g. "uses React"
#     becoming "prefers vanilla HTML/CSS/JS") instead of just accumulating.
#   - Both the SQLite row and the Chroma vector always store category +
#     importance together, so retrieval can rank by relevance AND importance
#     without needing a second round-trip to SQLite.

def store_memory_fact(user_id: int, text: str, category: str = "auto",
                       source_chat_id: int | None = None, dedup: bool = True,
                       importance: float = 0.7, confidence: float = 0.8):
    """Insert, update, or skip one atomic memory fact for a user.
    Returns the memory id that was written to (new or existing), or None if
    skipped as a true duplicate — OR rejected outright.

    Professor AI's memory is restricted by spec to academic-preference facts
    only (see MEMORY_CATEGORIES). Unlike a general assistant that falls back
    to a catch-all "other" bucket for anything that doesn't fit its taxonomy,
    here an unrecognized category is a hard rejection, not a relabel — a
    silent "other" fallback would have been exactly the loophole that let
    arbitrary text (including, worst case, a stray note about a student)
    into long-term storage via the manual /api/memory endpoint. Every caller,
    including manual entry, must pass one of the seven allowed categories."""
    text = (text or "").strip()
    if not text:
        return None
    if category not in MEMORY_CATEGORIES:
        return None
    importance = max(0.0, min(1.0, importance))
    confidence = max(0.0, min(1.0, confidence))

    col = get_user_memory_collection(user_id)

    if dedup:
        try:
            if col.count() > 0:
                nearest = col.query(query_embeddings=[_memory_embed(text, is_query=False)], n_results=1,
                                     include=["distances", "metadatas"])
                dists = (nearest.get("distances") or [[]])[0]
                ids = (nearest.get("ids") or [[]])[0]
                metas = (nearest.get("metadatas") or [[]])[0]
                if dists:
                    dist = dists[0]
                    if dist < MEMORY_DEDUP_DISTANCE:
                        # True duplicate — don't stack it, just reinforce importance
                        # if this occurrence carries a stronger signal than before.
                        near_id = int(ids[0].split("_", 1)[1])
                        existing_row = get_memory_row(user_id, near_id)
                        if existing_row and importance > existing_row["importance"]:
                            update_memory_row(user_id, near_id, existing_row["memory_text"],
                                               existing_row["category"], importance,
                                               max(confidence, existing_row["confidence"]))
                            _sync_memory_metadata(col, near_id, existing_row["category"], importance)
                        return None
                    if dist < MEMORY_UPDATE_DISTANCE:
                        near_category = (metas[0] or {}).get("category", category)
                        if near_category == category:
                            # Same topic, evolved content — overwrite in place.
                            near_id = int(ids[0].split("_", 1)[1])
                            update_memory_row(user_id, near_id, text, category, importance, confidence)
                            try:
                                col.update(ids=[f"mem_{near_id}"], documents=[text],
                                           metadatas=[{"category": category, "importance": importance}],
                                           embeddings=[_memory_embed(text, is_query=False)])
                            except Exception:
                                pass
                            return near_id
        except Exception:
            pass  # dedup/update is best-effort; a failed check should never block storage

    mem_id = insert_memory_row(user_id, text, category, source_chat_id, importance, confidence)
    try:
        col.add(ids=[f"mem_{mem_id}"], documents=[text],
                metadatas=[{"category": category, "importance": importance}],
                embeddings=[_memory_embed(text, is_query=False)])
    except Exception:
        # embedding failed — remove the orphaned SQLite row rather than leave
        # a fact that can never be retrieved
        delete_memory_row(user_id, mem_id)
        return None
    return mem_id


def _sync_memory_metadata(col, memory_id: int, category: str, importance: float):
    """Keep the Chroma metadata's importance score in step with SQLite after a
    reinforcement update, without re-embedding the (unchanged) document text."""
    try:
        col.update(ids=[f"mem_{memory_id}"], metadatas=[{"category": category, "importance": importance}])
    except Exception:
        pass


def delete_memory_fact(user_id: int, memory_id: int) -> bool:
    ok = delete_memory_row(user_id, memory_id)
    if ok:
        try:
            get_user_memory_collection(user_id).delete(ids=[f"mem_{memory_id}"])
        except Exception:
            pass
    return ok


def clear_user_memory(user_id: int):
    clear_user_memory_rows(user_id)
    drop_user_memory_collection(user_id)


# Defense-in-depth against student data ever entering long-term memory: even
# though the extraction prompt explicitly forbids it and manual entry requires
# a valid academic-preference category, a fact that slips through tagged with
# a technically-valid category but still referencing a specific student's
# grade/attendance/rank/etc. gets caught here before it's ever stored. Used
# by both run_memory_extraction (below) and the manual /api/memory endpoint.
_STUDENT_DATA_GUARD_RE = re.compile(
    r"\b(student|roll\s*no|marks?|grade[sd]?|score[sd]?|attendance|percentage|"
    r"cgpa|gpa|rank)\b", re.IGNORECASE
)

_TRIVIAL_MESSAGE_RE = re.compile(
    r"^(hi+|hey+|hello+|yo+|sup|ok(ay)?|k|thanks?|thank\s?you|thx|bye|goodbye|"
    r"good\s?(morning|night|evening|afternoon)|test(ing)?|hmm+)[\s!.,?]*$",
    re.IGNORECASE,
)


def _is_trivial_message(text: str) -> bool:
    """Greetings/acks/one-word filler are too short and generic for embedding
    distance to meaningfully judge — skip memory retrieval entirely rather
    than risk a spuriously 'close' hit surfacing an unrelated fact."""
    t = (text or "").strip()
    return len(t) < 4 or bool(_TRIVIAL_MESSAGE_RE.match(t))


def retrieve_memory(user_id: int, query: str, k: int = MEMORY_TOP_K) -> list[dict]:
    """Return up to k relevant remembered facts for this user, ranked by a
    blend of semantic closeness to the query AND how important the fact is
    (a low-importance tangential fact shouldn't outrank a high-importance
    core fact just because it's a hair closer in embedding space). Returns
    dicts with text + category so the prompt can group them for the model.
    Facts farther than MEMORY_MAX_DISTANCE from the query are dropped
    entirely — recall only kicks in when something is actually relevant,
    not on every single message regardless of topic."""
    if _is_trivial_message(query):
        return []
    try:
        col = get_user_memory_collection(user_id)
        count = col.count()
        if count == 0:
            return []
        pool = min(max(k * 3, MEMORY_CANDIDATE_POOL), count)
        results = col.query(query_embeddings=[_memory_embed(query, is_query=True)], n_results=pool,
                             include=["documents", "metadatas", "distances"])
        docs = results["documents"][0]
        metas = results["metadatas"][0]
        dists = results["distances"][0]
        if not docs:
            return []

        # TEMPORARY DIAGNOSTIC — prints the real distance so MEMORY_MAX_DISTANCE
        # can be set from actual data instead of guessed. Safe to delete once
        # recall behaves the way you want.
        best_dist = min(dists)
        print(f"[memory recall] query={query!r} best_dist={best_dist:.4f} "
              f"cutoff={MEMORY_MAX_DISTANCE} nearest={docs[dists.index(best_dist)]!r}")

        scored = []
        w = MEMORY_RETRIEVAL_IMPORTANCE_WEIGHT
        for doc, meta, dist in zip(docs, metas, dists):
            if dist > MEMORY_MAX_DISTANCE:
                continue  # not actually relevant to this message — leave it out
            closeness = 1.0 / (1.0 + dist)
            importance = float((meta or {}).get("importance", 0.7))
            score = (1 - w) * closeness + w * importance
            scored.append((score, doc, (meta or {}).get("category", "other")))

        scored.sort(key=lambda x: x[0], reverse=True)
        return [{"text": doc, "category": cat} for _, doc, cat in scored[:k]]
    except Exception:
        return []


MEMORY_EXTRACTION_PROMPT = """You extract durable academic-preference facts about a PROFESSOR from a slice of their conversation with Professor AI, a faculty teaching assistant.

You may ONLY extract facts that fall into these categories — nothing else, ever:

- "identity": the professor's own name
- "department": their department or faculty
- "subject": subjects/courses they teach
- "teaching_style": how they prefer material explained or structured (e.g. "prefers example-driven explanations", "likes analogies for abstract concepts")
- "question_pattern": how they like assessment questions built (e.g. "prefers scenario-based MCQs", "wants numerical problems with worked solutions")
- "output_format": their preferred output format/length/structure for generated content (e.g. "prefers DOCX over PDF", "likes concise bullet-point lecture notes")
- "frequent_prompt": a request pattern they clearly use repeatedly across sessions (e.g. "often asks for 20-question quizzes with an answer key")

You must NEVER extract, and must silently discard:
- Any personal conversation, opinion, or chit-chat that isn't one of the categories above
- Anything about a specific STUDENT — a name, a grade, a score, an evaluation, disciplinary information, or any other student-identifying or student-private detail. This holds even if the professor typed it themselves; it is never long-term memory material.
- Sensitive information of any kind (health, financial, personal life)
- The content of a one-off question or its answer, or anything about the subject matter being discussed rather than about the professor's own preferences
- Anything generic that wouldn't help personalize a FUTURE, unrelated session

For each fact, write it as a short, standalone, third-person sentence (e.g. "Professor teaches Data Structures and Algorithms.") and rate its importance from 0.0 to 1.0 (identity/department/subject are usually high; a one-off formatting preference lower).

Respond with ONLY a JSON array of objects, nothing else, in this exact shape:
[{{"memory": "Professor teaches Data Structures and Algorithms.", "category": "subject", "importance": 0.85}}]

If nothing in the excerpt fits one of the allowed categories, respond with exactly: []

Conversation excerpt:
{excerpt}
"""


def _parse_fact_objects(raw: str) -> list[dict]:
    """Parses the model's JSON output into normalized fact dicts, tolerant of
    a markdown code fence, a stray legacy string list, malformed categories,
    or out-of-range importance — any of which would otherwise silently drop
    real facts or crash the extraction thread.

    Professor AI has no catch-all category: a fact the model tags with
    anything outside MEMORY_CATEGORIES is DROPPED here, not relabeled and
    kept. That's a deliberate one-way filter — if the extraction model ever
    tries to sneak through something like a student's name or a stray
    personal remark under an invented category, it's discarded rather than
    quietly falling into a general "other" bucket that would defeat the
    whole point of restricting memory to academic preferences."""
    raw = (raw or "").strip()
    if raw.startswith("```"):
        raw = raw.strip("`")
        if raw.lower().startswith("json"):
            raw = raw[4:]
        raw = raw.strip()
    try:
        parsed = json.loads(raw)
    except Exception:
        return []
    if not isinstance(parsed, list):
        return []

    facts = []
    for item in parsed:
        if isinstance(item, str):
            continue  # no category info at all — can't verify it's in-taxonomy, so skip it
        elif isinstance(item, dict):
            text = str(item.get("memory") or item.get("text") or "").strip()
            category = str(item.get("category") or "").strip().lower()
            try:
                importance = float(item.get("importance", 0.6))
            except (TypeError, ValueError):
                importance = 0.6
        else:
            continue
        if not text or category not in MEMORY_CATEGORIES:
            continue
        importance = max(0.0, min(1.0, importance))
        if importance < MEMORY_MIN_IMPORTANCE:
            continue
        facts.append({"text": text, "category": category, "importance": importance})
    return facts


def run_memory_extraction(chat_id: int, user_id: int):
    """Runs in a background thread after a chat reply is saved. Pulls the new
    messages since this chat's watermark, asks the model to extract durable,
    categorized facts with importance scores, stores them (updating existing
    facts in place when appropriate — see store_memory_fact), and advances
    the watermark. Any failure here is logged and swallowed — memory
    extraction must never affect the live chat."""
    try:
        new_msgs = get_new_messages_since_watermark(chat_id)
        if not new_msgs:
            return

        excerpt = "\n".join(f"{m['role'].upper()}: {m['content']}" for m in new_msgs)
        excerpt = excerpt[:MEMORY_EXCERPT_CHAR_CAP]
        prompt = MEMORY_EXTRACTION_PROMPT.format(excerpt=excerpt)

        resp = ollama.chat(model=MEMORY_MODEL, messages=[{"role": "user", "content": prompt}], stream=False)
        facts = _parse_fact_objects(resp["message"]["content"])

        # Auto-extracted facts get a fixed, honest confidence — a 7B model
        # self-rating its own certainty is not reliable, so this isn't derived
        # from the model; it's just lower than the 1.0 given to facts the user
        # states directly through the manual "remember this" endpoint.
        for fact in facts:
            if _STUDENT_DATA_GUARD_RE.search(fact["text"]):
                continue  # never let student-referencing text into long-term memory
            store_memory_fact(user_id, fact["text"], category=fact["category"],
                               source_chat_id=chat_id, dedup=True,
                               importance=fact["importance"], confidence=0.5)

        set_extraction_watermark(chat_id, new_msgs[-1]["id"])
    except Exception as e:
        print(f"[memory extraction] chat {chat_id} failed: {e}")


def maybe_trigger_memory_extraction(chat_id: int, user_id: int):
    """Fire-and-forget: only starts a background extraction run once enough
    new messages have piled up, so we're not calling the LLM after every
    single turn."""
    try:
        if count_new_messages_since_watermark(chat_id) >= MEMORY_EXTRACTION_THRESHOLD:
            threading.Thread(target=run_memory_extraction, args=(chat_id, user_id), daemon=True).start()
    except Exception as e:
        print(f"[memory extraction] trigger check failed for chat {chat_id}: {e}")


# ═══════════════════════════════════════════════════════
# DOCUMENT GENERATION ENGINE (Phase 1)
# ═══════════════════════════════════════════════════════
# Prompt -> structured content -> DOCX/PDF. The standalone /api/generate-document
# endpoint below stays prompt-only (no chat_id is ever passed to it, so there's
# nothing to ground it in). The chat-triggered path further down does have a
# chat_id and pulls in that chat's uploaded-document chunks as source material —
# see gather_doc_gen_source_context() / build_doc_gen_prompt().

_JSON_FENCE_RE = re.compile(r"^```(?:json)?|```$", re.IGNORECASE)
_TRAILING_COMMA_RE = re.compile(r",(\s*[\]}])")


def _strip_json_fence(raw: str) -> str:
    """Model output is sometimes wrapped in a ```json ... ``` fence despite
    being told not to — strip it before parsing rather than let json.loads blow up."""
    raw = (raw or "").strip()
    raw = _JSON_FENCE_RE.sub("", raw).strip()
    return raw


_JSON_REPAIR_PROMPT = """The text below was supposed to be one valid JSON object but has a syntax error.

--- BROKEN JSON ---
{broken}
--- END BROKEN JSON ---

Parser error: {error}

Return ONLY the corrected, valid JSON object — same structure and same content, just fix the syntax (e.g. escape a stray double-quote inside a string value, remove a trailing comma, close an unterminated string). Do not add commentary, explanation, or a code fence. Do not change any of the actual wording, only the JSON punctuation around it.
"""


def _parse_llm_json(raw: str, repair_with_model: str | None = None) -> dict:
    """Parses a JSON object out of raw LLM output, tolerating the ways a
    local model most commonly breaks strict JSON when writing anything
    longer than a couple of words per field:
      1. A ```json fence wrapped around it despite being told not to.
      2. A trailing comma before a closing ] or } (the model second-guesses
         a list's length mid-generation and leaves one behind).
      3. Literal control characters inside string values — json.loads is
         strict about these by default even though they're harmless.
    If all of that still fails and repair_with_model is given, makes ONE
    extra call asking the model to fix its own output's syntax without
    touching the content — this is what actually recovers from the case
    that was crashing document generation: a stray unescaped quote inside a
    long generated paragraph prematurely closing a JSON string. Regex alone
    can't safely fix that (it can't tell content quotes from structural
    ones), but the model that wrote the text usually can.
    Raises the parse error if nothing works, so callers keep their existing
    fallback/error-reporting behaviour."""
    cleaned = _strip_json_fence(raw)
    try:
        return json.loads(cleaned, strict=False)
    except json.JSONDecodeError:
        try:
            return json.loads(_TRAILING_COMMA_RE.sub(r"\1", cleaned), strict=False)
        except json.JSONDecodeError as e2:
            if not repair_with_model:
                raise e2
            try:
                fix_prompt = _JSON_REPAIR_PROMPT.format(broken=cleaned[:6000], error=str(e2))
                resp = ollama.chat(model=repair_with_model,
                                    messages=[{"role": "user", "content": fix_prompt}], stream=False)
                fixed = _strip_json_fence(resp["message"]["content"])
                return json.loads(fixed, strict=False)
            except Exception:
                raise e2


DOC_REQUIREMENTS_PROMPT = """You analyze a professor's document request and extract structured requirements. Do not write the document itself here — only describe it.

Determine:
- "title": a clean, specific Title Case document name
- "type": one of "lecture_notes", "revision_notes", "assignment", "quiz", "mcq_set", "descriptive_questions", "lab_manual", "question_paper", "answer_key", "lesson_plan", "teaching_schedule", "rubric_evaluation", "literature_review", "report", "summary", "notes"
- "length": a short description of the requested size if mentioned (e.g. "10 pages", "20 questions", "short", "detailed"), else "standard"

Respond with ONLY a JSON object, nothing else, in this exact shape:
{{"title": "DBMS Unit 3 Lecture Notes", "type": "lecture_notes", "length": "standard"}}

Request:
{prompt}
"""


# extract_document_requirements() below asks the model to classify a
# free-text request into one of the DOC_REQUIREMENTS_PROMPT's type strings.
# A 7B model doesn't always reproduce that exact enum value verbatim — it
# might return "mcq" instead of "mcq_set", "multiple_choice_questions",
# "MCQs", etc. Previously that classification was used completely
# unvalidated: whatever string came back was passed straight through, and
# _DOC_TYPE_GUIDANCE.get(doc_type) would silently return None for anything
# that wasn't an exact key match — which meant the type-specific formatting
# instructions (e.g. "write actual (a)(b)(c)(d) options with the answer
# marked") never got injected, and generation quietly fell back to writing
# generic explanatory notes ABOUT the topic instead of the requested content
# type. This is what caused "generate me 5 MCQs" to come back as an
# Introduction/Key-Characteristics essay instead of five actual questions.
#
# Two layers of defense against that now:
#   1. _infer_doc_type_from_prompt: a deterministic keyword check against the
#      professor's OWN wording. This runs first and, if it matches, is
#      trusted over whatever the model classifies — the professor's own
#      words ("generate 5 MCQs on X") are a more reliable signal than an LLM
#      re-interpreting them.
#   2. _normalize_doc_type: a synonym map applied to whatever the model
#      returns, for the cases layer 1 doesn't confidently match.
_DOC_TYPE_KEYWORD_PATTERNS = [
    # (regex, canonical type) — checked in order, first match wins, so more
    # specific/less ambiguous phrases are listed before broader ones.
    (r"\banswer\s*key\b", "answer_key"),
    (r"\bquestion\s*paper\b", "question_paper"),
    (r"\bdescriptive\s*questions?\b", "descriptive_questions"),
    (r"\bmcqs?\b|\bmulti(?:ple)?[\s-]*choice\b", "mcq_set"),
    (r"\bquiz(?:zes)?\b", "quiz"),
    (r"\blab\s*manual\b", "lab_manual"),
    (r"\blesson\s*plan\b", "lesson_plan"),
    (r"\bteaching\s*schedule\b", "teaching_schedule"),
    (r"\brubric\b|\bevaluat(?:e|ion)\s*report\b", "rubric_evaluation"),
    (r"\bliterature\s*review\b", "literature_review"),
    (r"\blecture\s*notes?\b", "lecture_notes"),
    (r"\brevision\s*notes?\b", "revision_notes"),
    (r"\bassignment\b", "assignment"),
]

_DOC_TYPE_SYNONYMS = {
    "mcq": "mcq_set", "mcqs": "mcq_set", "multiple_choice": "mcq_set",
    "multiple_choice_questions": "mcq_set", "multiplechoicequestions": "mcq_set",
    "question_papers": "question_paper", "answer_keys": "answer_key",
    "lab_manuals": "lab_manual", "lesson_plans": "lesson_plan",
    "teaching_schedules": "teaching_schedule", "rubric": "rubric_evaluation",
    "evaluation": "rubric_evaluation", "evaluation_report": "rubric_evaluation",
    "literature_reviews": "literature_review", "lecture_note": "lecture_notes",
    "revision_note": "revision_notes", "assignments": "assignment",
    "descriptive_question": "descriptive_questions",
}


def _infer_doc_type_from_prompt(prompt: str) -> str | None:
    p_lower = (prompt or "").lower()
    for pattern, canonical in _DOC_TYPE_KEYWORD_PATTERNS:
        if re.search(pattern, p_lower):
            return canonical
    return None


def _normalize_doc_type(doc_type: str) -> str:
    doc_type = (doc_type or "").strip().lower().replace(" ", "_").replace("-", "_")
    return _DOC_TYPE_SYNONYMS.get(doc_type, doc_type)


# Document types where the professor's request implies a specific COUNT of
# discrete items (questions) rather than a page/section count. Used by
# generate_document_content() to force an exact total instead of letting the
# generic "N sections x 6-9 points each" depth rule multiply it out.
_QUESTION_STYLE_TYPES = {"mcq_set", "quiz", "descriptive_questions", "answer_key"}
_EXPLICIT_COUNT_RE = re.compile(r"\b(\d{1,3})\b")


def _extract_explicit_count(text: str) -> int | None:
    """Pulls a plausible explicit item count (e.g. the "5" in "5 MCQs" or
    "20" in "20-question quiz") out of free text. Rejects absurd values
    (0 or >200) as more likely to be something else entirely (a year, a
    course code) than a genuine question count."""
    m = _EXPLICIT_COUNT_RE.search(text or "")
    if not m:
        return None
    n = int(m.group(1))
    return n if 0 < n <= 200 else None


def extract_document_requirements(prompt: str) -> dict:
    """Step 2 of the workflow: turns a free-text request into {title, type, length}.
    Falls back to a reasonable default derived from the prompt itself if the
    model's output can't be parsed, so a generation request never hard-fails
    on this step alone."""
    keyword_type = _infer_doc_type_from_prompt(prompt)
    try:
        resp = ollama.chat(model=DOC_GEN_MODEL,
                            messages=[{"role": "user", "content": DOC_REQUIREMENTS_PROMPT.format(prompt=prompt)}],
                            stream=False)
        data = _parse_llm_json(resp["message"]["content"], repair_with_model=DOC_GEN_MODEL)
        title = str(data.get("title") or "").strip() or prompt.strip()[:60].title()
        # The professor's own wording (layer 1) wins over the model's
        # classification whenever it matched something; only fall through to
        # the model's answer — normalized (layer 2) — when it didn't.
        doc_type = keyword_type or _normalize_doc_type(str(data.get("type") or "notes"))
        length = str(data.get("length") or "standard").strip()
        # TEMPORARY DIAGNOSTIC — shows exactly what type/length this request
        # resolved to, so a wrong-format result can be traced to "wrong type
        # detected" vs. "type was right but the model ignored the guidance"
        # instead of guessing. Safe to delete once generation looks right.
        print(f"[doc-gen requirements] prompt={prompt!r} -> keyword_type={keyword_type!r} "
              f"llm_type={data.get('type')!r} -> resolved doc_type={doc_type!r} length={length!r}")
        return {"title": title or "Untitled Document", "type": doc_type, "length": length}
    except Exception:
        return {"title": (prompt.strip()[:60] or "Untitled Document").title(),
                "type": keyword_type or "notes", "length": "standard"}


def _estimate_section_count(length: str) -> int:
    """Rough heuristic mapping a requested length ('10 pages', 'short', 'detailed')
    to a target section count, so a 10-page request doesn't come back as 3
    thin sections or a 200-section wall of text."""
    length = (length or "").lower()
    m = re.search(r"(\d+)", length)
    if m:
        return max(3, min(DOC_MAX_SECTIONS, int(m.group(1))))
    if "short" in length or "brief" in length:
        return 5
    if "detailed" in length or "long" in length or "comprehensive" in length:
        return 12
    return 8  # "standard" / unspecified — biased toward substantial rather than skeletal


# Type-specific guidance folded into the generation prompt below. Every type
# still emits the same universal {heading, content|points|table} JSON shape
# (so render_docx/render_pdf need no changes at all), but each gets extra
# instructions on how to use those three formats to represent that
# document's actual structure — an MCQ set needs answers marked, an answer
# key needs question+model-answer+marks, a rubric needs scored criteria.
_DOC_TYPE_GUIDANCE = {
    "mcq_set": (
        "This is a multiple-choice question set. Use \"points\" sections. Each point is one complete "
        "question formatted as: \"Q1. <question text> (a) <option> (b) <option> (c) <option> (d) <option> "
        "— Answer: <letter>\". Group into sections by topic if the request covers several topics."
    ),
    "quiz": (
        "This is a quiz. Mix question types as appropriate (MCQ, short-answer, true/false). For MCQs use "
        "the format \"Q1. <question> (a)...(b)...(c)...(d)... — Answer: <letter>\" in a \"points\" list. "
        "For short-answer questions, list the question and give the expected answer in parentheses."
    ),
    "descriptive_questions": (
        "This is a set of descriptive/long-answer questions. Use \"points\" sections; each point is one "
        "complete question, optionally with a suggested mark allocation in parentheses, e.g. "
        "\"Explain the CAP theorem with examples. (10 marks)\". Do not include model answers."
    ),
    "answer_key": (
        "This is an answer key. Use \"points\" sections. Each point covers one question as "
        "\"Q1. <model answer, complete and exam-ready> (<marks> marks)\". Be thorough and precise — "
        "this is what will be used to grade student submissions."
    ),
    "question_paper": (
        "This is a formal question paper. Structure sections by part (e.g. \"Section A — 2 Marks Each\", "
        "\"Section B — 10 Marks Each\") using \"points\" lists of numbered questions with marks noted. "
        "Include a brief instructions section at the top (duration, total marks, instructions to candidates) "
        "as a \"content\" section."
    ),
    "lab_manual": (
        "This is a lab manual. For each experiment, use a section with \"content\" covering the Aim and "
        "Theory, followed by a \"points\" section for Procedure (numbered steps), and where relevant "
        "another \"points\" or \"content\" section for expected output / sample results. Include any "
        "example code in the content text using fenced-style formatting."
    ),
    "lesson_plan": (
        "This is a lesson plan. Use sections for Learning Objectives (points), Materials Needed (points), "
        "Lesson Flow / Timing (table with columns like Time, Activity, Description), and Assessment "
        "(content or points)."
    ),
    "teaching_schedule": (
        "This is a weekly teaching schedule. Use a \"table\" section per period covered, with columns like "
        "Week, Topic, Learning Outcome, Suggested Activity — or one table for the whole schedule if it fits."
    ),
    "rubric_evaluation": (
        "This is a rubric-based evaluation report. Use a \"table\" section with columns "
        "[Criterion, Weight/Max Marks, Score Awarded, Feedback] — one row per rubric criterion — followed "
        "by a \"content\" section with overall constructive feedback and concrete suggestions for improvement. "
        "Base every score and comment strictly on the submission/rubric text provided; never invent a score "
        "for something not actually assessable from what was given."
    ),
    "lecture_notes": (
        "These are lecture notes for classroom delivery. Favor clear explanations with concrete examples "
        "and analogies for difficult concepts, organized in the order a lecture would flow."
    ),
}


DOC_CONTENT_PROMPT = """You are Professor AI, generating a structured educational document for a faculty member's classroom or assessment use.

Document title: {title}
Document type: {doc_type}
Requested length: {length}
Original request: {prompt}
{type_guidance}
Write clear, accurate, well-organized, professional educational content broken into logical sections. Aim for {section_target} sections covering the topic thoroughly. This needs to be a genuinely substantial, classroom-ready document — not a skeleton outline — so follow these depth requirements for every section unless the type-specific guidance above overrides it:
- A "content" paragraph must be at least 150-250 words: real explanation with definitions, reasoning, and examples where relevant, not a one- or two-sentence summary.
- A "points" list must have at least 6-9 substantive entries (fewer only if the type guidance above calls for a specific smaller count, e.g. a fixed number of quiz questions), each a full, usable entry, not a short fragment or bare keyword.
- A "table" must have at least 4-5 data rows beyond the header — enough to be genuinely useful, not a token example.
Depth and coverage are the priority. Do not pad with filler, restate the heading as a sentence, or repeat the same point across sections. If the topic doesn't naturally support {section_target} distinct sections, cover it as thoroughly as it warrants with fewer, deeper sections rather than splitting it into thin, repetitive ones.

For each section, choose the format that reads best:
- "content": a prose paragraph string — for explanations, definitions, narrative content.
- "points": an array of short strings — for lists, steps, questions, comparisons, enumerable items.
- "table": a list of rows (each row a list of cell strings), first row is the header — only for naturally tabular content (comparisons, schedules, rubrics, specifications).

A section has exactly ONE of "content", "points", or "table".

Output strict, valid JSON: escape every literal double-quote and backslash that appears inside a string value (e.g. write \" not "), and do not leave a trailing comma before a closing ] or }}.

Respond with ONLY a JSON object, nothing else, in this exact shape:
{{
  "title": "{title}",
  "sections": [
    {{"heading": "Introduction", "content": "..."}},
    {{"heading": "Key Points", "points": ["...", "..."]}},
    {{"heading": "Comparison", "table": [["Algorithm", "Preemptive"], ["FCFS", "No"]]}}
  ]
}}

Do not include a References, Citations, or Bibliography section unless the document type is literature_review.{grounding_rule}
"""

# Swapped in when the request came with actual retrieved source material
# (see gather_doc_gen_source_context) — replaces the generic "don't fabricate
# statistics" line with a harder rule, since here there IS a ground truth to
# be unfaithful to, not just an abstract caution.
_GROUNDED_RULE = (" The \"Original request\" below begins with source material pulled "
                   "from the user's own uploaded document(s). Base every section on that "
                   "source material — do not introduce facts, figures, or points that "
                   "aren't supported by it. If the source material doesn't cover something "
                   "the request asks for, leave it out rather than inventing it.")
_UNGROUNDED_RULE = " Do not fabricate specific statistics or sources."


def generate_document_content(title: str, doc_type: str, length: str, prompt: str,
                               grounded: bool = False) -> dict:
    """Step 3 of the workflow: the model generates STRUCTURE (sections with
    headings + content/points/table), not a file — rendering into an actual
    document happens separately in render_docx/render_pdf. Raises on failure
    so the route can return a clear error instead of writing an empty file.

    grounded=True means `prompt` already has retrieved source-document
    excerpts prepended to it (see gather_doc_gen_source_context) — in that
    case the model is instructed to stick to that material instead of the
    generic "don't fabricate" caution used for pure invent-from-scratch requests."""
    doc_type_norm = doc_type.lower()
    grounding_rule = _GROUNDED_RULE if grounded else _UNGROUNDED_RULE
    guidance = _DOC_TYPE_GUIDANCE.get(doc_type_norm)

    # Question-style documents (MCQ sets, quizzes, descriptive question sets,
    # answer keys) need an EXACT total count, not "aim for N sections" — the
    # generic section-count heuristic below treats N as a section count and
    # the depth rule pads each section's points list to 6-9 entries, so
    # "generate 5 MCQs" was silently becoming 5 sections x 6-9 questions
    # each (30-45 questions) instead of 5. When an explicit count is present,
    # force a single section with exactly that many points instead.
    explicit_count = None
    if doc_type_norm in _QUESTION_STYLE_TYPES:
        explicit_count = _extract_explicit_count(length) or _extract_explicit_count(prompt)

    # TEMPORARY DIAGNOSTIC — safe to delete once generation looks right.
    print(f"[doc-gen content] doc_type_norm={doc_type_norm!r} "
          f"guidance_found={guidance is not None} explicit_count={explicit_count!r}")

    if explicit_count:
        section_target = 1
        count_rule = (
            f"\nThis document must contain EXACTLY {explicit_count} questions in total — no more, "
            f"no fewer. Put them all in ONE \"points\" section (do not split into multiple sections, "
            f"and do not pad or trim to hit any other target count).\n"
        )
        type_guidance = f"\nDocument-type-specific instructions: {guidance}{count_rule}" if guidance else count_rule
    else:
        section_target = _estimate_section_count(length)
        type_guidance = f"\nDocument-type-specific instructions: {guidance}\n" if guidance else ""

    filled_prompt = DOC_CONTENT_PROMPT.format(title=title, doc_type=doc_type, length=length,
                                               prompt=prompt, section_target=section_target,
                                               grounding_rule=grounding_rule, type_guidance=type_guidance)
    resp = ollama.chat(model=DOC_GEN_MODEL, messages=[{"role": "user", "content": filled_prompt}],
                        stream=False, options={"num_predict": DOC_GEN_MAX_TOKENS})
    data = _parse_llm_json(resp["message"]["content"], repair_with_model=DOC_GEN_MODEL)

    clean_sections = []
    for s in (data.get("sections") or [])[:DOC_MAX_SECTIONS]:
        if not isinstance(s, dict):
            continue
        heading = str(s.get("heading") or "").strip()
        if not heading:
            continue
        entry = {"heading": heading}
        if isinstance(s.get("points"), list) and s["points"]:
            entry["points"] = [str(p).strip() for p in s["points"] if str(p).strip()]
        elif isinstance(s.get("table"), list) and len(s["table"]) >= 2:
            rows = [[str(cell).strip() for cell in row] for row in s["table"] if isinstance(row, list) and row]
            if len(rows) >= 2:
                entry["table"] = rows
            else:
                entry["content"] = str(s.get("content") or "").strip()
        else:
            entry["content"] = str(s.get("content") or "").strip()
        if entry.get("points") or entry.get("table") or entry.get("content"):
            clean_sections.append(entry)

    if not clean_sections:
        raise ValueError("model returned no usable sections")

    return {"title": str(data.get("title") or title).strip() or title, "sections": clean_sections}


def make_safe_filename(title: str) -> str:
    """Collapses a title into a filesystem-safe base name (no extension),
    e.g. 'Revision Notes' -> 'Revision_Notes'."""
    base = re.sub(r"[^A-Za-z0-9]+", "_", title or "").strip("_")
    return (base or "Document")[:80]


def get_user_doc_dir(user_id: int) -> Path:
    d = GENERATED_DOCS_DIR / f"user_{user_id}"
    d.mkdir(parents=True, exist_ok=True)
    return d


def render_docx(structured: dict, out_path: Path):
    """DOCX rendering (Step 4): headings, paragraphs, bullet lists, and tables
    via python-docx, matching the doc's DOCX Generator responsibilities."""
    if not DOCX_OK:
        raise RuntimeError("python-docx is not installed. Run: pip install python-docx")

    doc = DocxDocument()
    doc.add_heading(structured["title"], level=0)

    for section in structured["sections"]:
        doc.add_heading(section["heading"], level=1)
        if "points" in section:
            for point in section["points"]:
                doc.add_paragraph(point, style="List Bullet")
        elif "table" in section:
            rows = section["table"]
            table = doc.add_table(rows=1, cols=len(rows[0]))
            try:
                table.style = "Light Grid Accent 1"
            except Exception:
                pass  # falls back to the default table style if this template lacks it
            for i, cell_text in enumerate(rows[0]):
                cell = table.rows[0].cells[i]
                cell.text = cell_text
                for p in cell.paragraphs:
                    for run in p.runs:
                        run.bold = True
            for row in rows[1:]:
                cells = table.add_row().cells
                for i, cell_text in enumerate(row):
                    if i < len(cells):
                        cells[i].text = cell_text
        else:
            doc.add_paragraph(section.get("content", ""))

    doc.save(str(out_path))


def _pdf_escape(text: str) -> str:
    """reportlab Paragraph content is mini-XML — escape special characters so
    section text containing '&', '<', or '>' doesn't break rendering."""
    return (text or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def render_pdf(structured: dict, out_path: Path):
    """PDF rendering (Step 4) via ReportLab, matching the doc's PDF Generator
    responsibilities: page layout, headings, paragraphs, and now tables too."""
    if not PDF_OK:
        raise RuntimeError("reportlab is not installed. Run: pip install reportlab")

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("DocTitle", parent=styles["Title"], fontSize=20, spaceAfter=20)
    heading_style = ParagraphStyle("DocHeading", parent=styles["Heading1"], fontSize=14,
                                    spaceBefore=14, spaceAfter=8, textColor=colors.HexColor("#1a1a2e"))
    body_style = ParagraphStyle("DocBody", parent=styles["BodyText"], fontSize=11, leading=16, spaceAfter=8)
    bullet_style = ParagraphStyle("DocBullet", parent=styles["BodyText"], fontSize=11,
                                   leading=15, leftIndent=6, spaceAfter=4)

    story = [Paragraph(_pdf_escape(structured["title"]), title_style), Spacer(1, 12)]

    for section in structured["sections"]:
        story.append(Paragraph(_pdf_escape(section["heading"]), heading_style))
        if "points" in section:
            items = [ListItem(Paragraph(_pdf_escape(p), bullet_style), leftIndent=18) for p in section["points"]]
            story.append(ListFlowable(items, bulletType="bullet", start="circle"))
        elif "table" in section:
            tbl = Table(section["table"], hAlign="LEFT")
            tbl.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2f3b52")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f2f2f2")]),
            ]))
            story.append(tbl)
        else:
            story.append(Paragraph(_pdf_escape(section.get("content", "")), body_style))
        story.append(Spacer(1, 6))

    pdf = SimpleDocTemplate(str(out_path), pagesize=A4,
                             leftMargin=0.9 * inch, rightMargin=0.9 * inch,
                             topMargin=0.9 * inch, bottomMargin=0.9 * inch)
    pdf.build(story)


def record_generated_document(user_id: int, title: str, doc_type: str, filename: str, fmt: str):
    conn = get_db()
    conn.execute(
        "INSERT INTO generated_documents (user_id, title, doc_type, filename, format) VALUES (?, ?, ?, ?, ?)",
        (user_id, title, doc_type, filename, fmt)
    )
    conn.commit()
    conn.close()


def list_generated_documents(user_id: int):
    conn = get_db()
    rows = conn.execute(
        "SELECT id, title, doc_type, filename, format, created_at FROM generated_documents "
        "WHERE user_id = ? ORDER BY created_at DESC",
        (user_id,)
    ).fetchall()
    conn.close()
    docs = [dict(r) for r in rows]
    for d in docs:
        d["download_url"] = f"/download/{d['filename']}"
    return docs
# ── chat-triggered document generation ──
# Lets a user type "generate me a word document of..." directly in the chat
# box and get a REAL file + working link, instead of the model hallucinating
# one. This is a cheap regex pre-check (fast, no extra LLM call on every
# message) that only fires the real /api/generate-document pipeline when it
# matches — everything else goes through the normal chat flow untouched.
_DOC_REQUEST_RE = re.compile(
    r"\b(generate|create|make|prepare|write|give me|produce|export|download)\b"
    r"[^.\n]{0,60}\b(word document|word doc|docx|pdf file|pdf|document|"
    r"lecture notes?|revision notes?|assignment|quiz|mcqs?|multiple[- ]choice|"
    r"descriptive questions?|lab manual|question paper|answer key|"
    r"lesson plan|teaching schedule|rubric|evaluation report|literature review)\b",
    re.IGNORECASE
)


def looks_like_document_request(text: str) -> tuple[bool, str]:
    """Cheap heuristic pre-check. Returns (is_doc_request, requested_format)."""
    t = (text or "").lower()
    if not _DOC_REQUEST_RE.search(t):
        return False, "pdf"
    if "docx" in t or "word doc" in t or "word document" in t:
        return True, "docx"
    if "pdf" in t:
        return True, "pdf"
    return True, "pdf"  # generic "document" mention defaults to pdf


def gather_doc_gen_source_context(chat_id: int, query: str) -> tuple[str, str | None]:
    """Pulls relevant chunks from THIS chat's uploaded documents (the same
    Chroma collection normal chat RAG uses) so document generation is
    grounded in what was actually uploaded.

    Before this, a request like "generate a summary of the important points
    from the PDF" reached generate_document_content() with nothing but that
    sentence — the LLM had no PDF content to summarize, so it filled the
    document with plausible-sounding but invented points instead. This makes
    the retrieved excerpts part of the prompt, the same way the normal chat
    endpoint already does for regular Q&A.

    Returns (context_block, source_filename). source_filename is whichever
    uploaded file contributed the most retrieved chunks — used to name the
    output file after the file it's actually summarizing — or None if this
    chat has no indexed documents, or nothing in them matched the query.
    """
    try:
        collection = get_chat_collection(chat_id)
        if collection.count() == 0:
            return "", None
    except Exception:
        return "", None

    hits = retrieve(collection, query, k=DOC_GEN_RAG_K)
    if not hits:
        return "", None

    filenames = [h["filename"] for h in hits if h.get("filename") and h["filename"] != "?"]
    source_filename = Counter(filenames).most_common(1)[0][0] if filenames else None

    context = "\n\n".join(f"[Source: {h['filename']}, page {h['page']}]\n{h['text']}" for h in hits)
    return context, source_filename


def build_doc_gen_prompt(chat_id: int, query: str) -> tuple[str, str | None, bool]:
    """Builds the prompt the generation pipeline actually works from.

    Priority order:
      1. This chat's uploaded documents — if any chunk is relevant to the
         request, ground the document in it. This is what makes "summarize
         the PDF" produce the PDF's actual points.
      2. The most recent assistant reply, for requests like "prepare a word
         document of the above info" that reference earlier chat content
         rather than an uploaded file.
      3. The bare request text, unchanged, if neither applies.

    Returns (prompt, source_filename, grounded) — source_filename names the
    uploaded file to base the output filename on (None if not applicable),
    and grounded tells generate_document_content() whether real source
    material is present so it should refuse to invent facts beyond it.
    """
    context, source_filename = gather_doc_gen_source_context(chat_id, query)
    if context:
        prompt = (
            f"{context}\n\n"
            f"User's request: {query}"
        )
        return prompt, source_filename, True

    if len(query) <= 200 and re.search(r"\b(above|previous|earlier|that|this)\b", query.lower()):
        recent = get_history(chat_id, limit=6)
        for m in reversed(recent):
            if m["role"] == "assistant":
                prompt = (
                    f"Base the document on this earlier conversation content:\n"
                    f"{m['content'][:4000]}\n\n"
                    f"User's request: {query}"
                )
                return prompt, None, False

    return query, None, False


def handle_chat_document_generation(chat_id: int, user_id: int, query: str, fmt: str):
    """Runs the real Step 2→4 pipeline synchronously and streams back a
    message with a genuine, working /download/... link, using the exact same
    SSE event shape the frontend already renders (token + done), so no
    frontend changes are needed."""
    def generate():
        try:
            gen_prompt, source_filename, grounded = build_doc_gen_prompt(chat_id, query)
            # extract_document_requirements only sees the (possibly context-prefixed)
            # prompt to pick a title/type/length — that's fine, it's not where content
            # gets invented. The actual content step below is where "grounded" matters.
            requirements = extract_document_requirements(gen_prompt)
            structured = generate_document_content(
                requirements["title"], requirements["type"], requirements["length"],
                gen_prompt, grounded=grounded
            )

            fmt_actual = fmt
            if fmt_actual == "docx" and not DOCX_OK:
                fmt_actual = "pdf"
            elif fmt_actual == "pdf" and not PDF_OK:
                fmt_actual = "docx"

            doc_dir = get_user_doc_dir(user_id)
            # Name the output after the uploaded file it's actually summarizing,
            # not the LLM's own made-up title — "Report.pdf" in, "Report.docx" out.
            # Falls back to the generated title only when there's no source file
            # (e.g. a from-scratch request, or one that referenced earlier chat text).
            if source_filename:
                base_name = make_safe_filename(Path(source_filename).stem)
            else:
                base_name = make_safe_filename(structured["title"])
            filename = f"{base_name}.{fmt_actual}"
            # Avoid silently clobbering an earlier generated file with the same
            # source name (e.g. re-running "summarize this PDF" twice) — append
            # a numeric suffix instead of overwriting the previous download.
            counter = 2
            while (doc_dir / filename).exists():
                filename = f"{base_name}_{counter}.{fmt_actual}"
                counter += 1
            if fmt_actual == "docx":
                render_docx(structured, doc_dir / filename)
            else:
                render_pdf(structured, doc_dir / filename)

            record_generated_document(user_id, structured["title"], requirements["type"], filename, fmt_actual)
            download_url = f"/download/{filename}"

            reply = f"I've generated **{structured['title']}**.\n\n[Download {filename}]({download_url})"

            yield f"data: {json.dumps({'type': 'token', 'content': reply})}\n\n"

            add_message(chat_id, "user", query)
            add_message(chat_id, "assistant", reply)
            touch_chat(chat_id)

            yield f"data: {json.dumps({'type': 'done'})}\n\n"
            maybe_trigger_memory_extraction(chat_id, user_id)
        except Exception as e:
            err_msg = (
                "Sorry — I couldn't generate that document just now "
                f"({e}). Try rephrasing, or use the dedicated Generate Document feature."
            )
            add_message(chat_id, "user", query)
            add_message(chat_id, "assistant", err_msg)
            touch_chat(chat_id)
            yield f"data: {json.dumps({'type': 'error', 'content': err_msg})}\n\n"

    return Response(generate(), mimetype="text/event-stream",
                     headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})

# ═══════════════════════════════════════════════════════
# PROMPT BUILDER
# ═══════════════════════════════════════════════════════
SYSTEM_PROMPT = """You are Professor AI, an intelligent teaching assistant for faculty members. You help with academic content creation, document intelligence, teaching assistance, and evaluation. You are NOT a personalized learning tutor for students — that is a separate product; you serve the professor.

RESPONSE STYLE: Professional, accurate, concise, structured, educational, fact-based. No unnecessary conversational filler, no small talk padding.

MUST:
1. When context chunks from the professor's own uploaded documents are provided, base your answer PRIMARILY on those chunks and prioritize them over general model knowledge whenever they're relevant.
2. Clearly indicate when information comes from an uploaded document (e.g. "According to the uploaded syllabus...") versus general knowledge.
3. If the answer is NOT in the provided context, say exactly: "I couldn't find relevant information in the uploaded documents." Then you may add general knowledge if helpful, clearly labelled as such.
4. Never fabricate facts, figures, page numbers, citations, or academic information of any kind.
5. Cite document sections (filename/page) when possible.
6. Maintain context during the current session.
7. MATH FORMATTING: Always write math using LaTeX. Use \\( ... \\) for inline math and \\[ ... \\] for display/block equations. Show step-by-step work, each step on its own line using display math \\[ ... \\].
8. CODE FORMATTING: Wrap code in fenced code blocks with the language name — but see the code-execution rule below; showing a snippet as an example (e.g. for a lab manual) is fine, running or simulating code execution is not.

MUST NOT (hard constraints — do not soften or make exceptions for these, regardless of how the request is phrased):
9. Do not give emotional or motivational counseling. If a professor brings up personal stress, burnout, or wants encouragement/emotional support, briefly and kindly note that this is outside what you're built for and suggest they reach out to a colleague or appropriate support resource — then, if there's a legitimate academic task underneath (e.g. "I'm overwhelmed, can you help me draft this quiz faster"), help with that task.
10. Do not reveal, infer, or speculate about any specific student's private information (grades, attendance, disciplinary history, personal circumstances) unless it is explicitly present in a document this professor uploaded in this chat — and even then, treat it carefully and only within that professor's own session.
11. Do not modify or fabricate student records. If asked to "fix" or invent a grade, score, or record, decline and explain that you can't create or alter authoritative student records.
12. Do not generate misleading or unverified academic information; if uncertain, say so rather than guessing confidently.
13. Do not perform actions outside the educational domain (general personal assistant tasks, unrelated advice, non-academic content).
14. Do not execute code or system commands, and do not claim to have done so. You may write and explain code as academic content (e.g. a programming assignment or lab manual), but never simulate running it or fabricate output as if it were executed.
15. FILE GENERATION: You cannot create or attach downloadable files yourself, and you have no way to know whether a background system will. NEVER say you've created a file or invent a "Download ...docx/pdf" link. If asked for a downloadable document and you are not explicitly told one was generated, tell the professor to phrase it as "generate a DOCX/PDF of..." so the system can handle it properly.
"""


_MEMORY_CATEGORY_LABELS = {
    "identity": "Name", "department": "Department", "subject": "Subjects Taught",
    "teaching_style": "Teaching Style", "question_pattern": "Question Pattern",
    "output_format": "Output Format", "frequent_prompt": "Frequently Used Prompts",
}


def _format_memory_block(memories: list) -> str:
    """Turns a list of {"text", "category"} memory dicts into the
    'PROFESSOR'S ACADEMIC PROFILE' block appended to a system prompt.
    Shared by build_messages (text chat) and the vision chat-image endpoint
    so both surfaces present memory the same way."""
    grouped: dict[str, list[str]] = {}
    for m in memories:
        if isinstance(m, dict):
            text, cat = m.get("text", ""), m.get("category", "other")
        else:
            text, cat = str(m), "other"  # tolerate the old plain-string shape
        if text:
            grouped.setdefault(cat, []).append(text)

    lines = []
    for cat, texts in grouped.items():
        label = _MEMORY_CATEGORY_LABELS.get(cat, cat.title())
        for t in texts:
            lines.append(f"- [{label}] {t}")
    mem_block = "\n".join(lines)

    return (
        "\n\nPROFESSOR'S ACADEMIC PROFILE (remembered from earlier sessions):\n"
        f"{mem_block}\n\n"
        "Use this naturally where it's relevant to personalize your answer (e.g. matching their "
        "preferred question pattern or output format). Don't force it into every response, and "
        "don't announce that you're 'recalling memory'."
    )


def build_messages(history: list, query: str, chunks: list, memories: list | None = None) -> list:
    system_content = SYSTEM_PROMPT
    if memories:
        system_content += _format_memory_block(memories)
    messages = [{"role": "system", "content": system_content}]
    messages.extend(history)
    if chunks:
        ctx_parts = [
            f"[Chunk {i} | File: {c['filename']} | Page: {c['page']}]\n{c['text']}"
            for i, c in enumerate(chunks, 1)
        ]
        context_block = "\n\n---\n\n".join(ctx_parts)
        user_content = (
            f"Use the following document excerpts to answer the question.\n\n"
            f"=== DOCUMENT CONTEXT ===\n{context_block}\n\n"
            f"=== QUESTION ===\n{query}"
        )
    else:
        user_content = query
    messages.append({"role": "user", "content": user_content})
    return messages


# ═══════════════════════════════════════════════════════
# ROUTES — AUTH PAGES
# ═══════════════════════════════════════════════════════
@app.route("/login")
def login_page():
    if "username" in session:
        return redirect(url_for("index"))
    return render_template("login.html")


@app.route("/api/auth/signup", methods=["POST"])
def signup():
    data = request.get_json(force=True)
    username = (data.get("username") or "").strip().lower()
    name = (data.get("name") or "").strip()
    password = data.get("password") or ""
    department = (data.get("department") or "").strip()
    subjects = (data.get("subjects") or "").strip()
    if not username or not name or not password:
        return jsonify({"error": "All fields are required."}), 400
    if len(password) < 4:
        return jsonify({"error": "Password must be at least 4 characters."}), 400
    if len(username) < 3:
        return jsonify({"error": "Username must be at least 3 characters."}), 400
    if not create_user(username, name, password, role="professor"):
        return jsonify({"error": "That username is already taken."}), 409
    session["username"] = username
    session["name"] = name
    session["role"] = "professor"

    # Seed the account's academic-preference memory immediately from what
    # they gave us at signup, rather than waiting for these to come up
    # naturally in conversation and get picked up by extraction — a name
    # and department stated once at signup shouldn't need to be repeated in
    # chat before Professor AI "knows" it. Uses the same restricted
    # MEMORY_CATEGORIES as everywhere else; nothing here is a personal
    # conversation or sensitive data, just the academic profile fields the
    # spec explicitly says are fine to remember.
    user = get_user(username)
    if user:
        store_memory_fact(user["id"], f"Professor's name is {name}.", category="identity",
                           dedup=False, importance=0.9, confidence=1.0)
        if department:
            store_memory_fact(user["id"], f"Department: {department}.", category="department",
                               dedup=False, importance=0.9, confidence=1.0)
        if subjects:
            store_memory_fact(user["id"], f"Subjects taught: {subjects}.", category="subject",
                               dedup=False, importance=0.9, confidence=1.0)

    return jsonify({"success": True, "name": name})


@app.route("/api/auth/login", methods=["POST"])
def login():
    data = request.get_json(force=True)
    username = (data.get("username") or "").strip().lower()
    password = data.get("password") or ""
    user = get_user(username)
    if not user or not check_password_hash(user["password_hash"], password):
        return jsonify({"error": "Incorrect username or password."}), 401
    session["username"] = user["username"]
    session["name"] = user["name"]
    session["role"] = user["role"]
    return jsonify({"success": True, "name": user["name"]})


@app.route("/api/auth/logout", methods=["POST"])
def logout():
    session.clear()
    return jsonify({"success": True})


@app.route("/api/auth/me", methods=["GET"])
def me():
    if "username" not in session:
        return jsonify({"authenticated": False}), 200
    return jsonify({
        "authenticated": True,
        "username": session["username"],
        "name": session.get("name", session["username"]),
        "role": session.get("role", "professor"),
    })


# ═══════════════════════════════════════════════════════
# ROUTES — CHATS (create / list / rename / delete)
# ═══════════════════════════════════════════════════════
@app.route("/api/chats", methods=["GET"])
@login_required
def api_list_chats():
    user = get_user(session["username"])
    return jsonify({"chats": list_chats(user["id"])})


@app.route("/api/chats", methods=["POST"])
@login_required
def api_create_chat():
    user = get_user(session["username"])
    data = request.get_json(silent=True) or {}
    title = (data.get("title") or "New chat").strip()[:100]
    chat_id = create_chat(user["id"], title)
    return jsonify({"success": True, "chat_id": chat_id, "title": title})


@app.route("/api/chats/<chat_id>", methods=["PATCH"])
@login_required
@chat_access_required
def api_rename_chat(chat_id):
    data = request.get_json(force=True)
    title = (data.get("title") or "").strip()[:100]
    if not title:
        return jsonify({"error": "Title required"}), 400
    rename_chat(chat_id, title)
    return jsonify({"success": True})


@app.route("/api/chats/<chat_id>", methods=["DELETE"])
@login_required
@chat_access_required
def api_delete_chat(chat_id):
    # wipe the chat's documents from the vector store, then the DB row
    # (messages + chat_files cascade via ON DELETE CASCADE)
    drop_chat_collection(chat_id)
    for f in list_chat_files(chat_id):
        (UPLOAD_DIR / f["stored_name"]).unlink(missing_ok=True)
    delete_chat_row(chat_id)
    return jsonify({"success": True})


# ═══════════════════════════════════════════════════════
# ROUTES — UPLOAD (scoped to one chat)
# ═══════════════════════════════════════════════════════
@app.route("/api/chats/<chat_id>/upload", methods=["POST"])
@login_required
@chat_access_required
def upload_file(chat_id):
    file = request.files.get("file")
    if not file or not file.filename:
        return jsonify({"error": "No file provided"}), 400

    ext = Path(file.filename).suffix.lower()
    if ext not in [".pdf", ".txt", ".docx", ".pptx"]:
        return jsonify({"error": f"Unsupported format '{ext}'. Use PDF, TXT, DOCX, or PPTX."}), 400

    # stored name is namespaced by chat_id so the same filename in two
    # chats never overwrites the other chat's file on disk
    stored_name = f"chat{chat_id}_{file.filename}".replace("/", "_")
    save_path = UPLOAD_DIR / stored_name
    file.save(str(save_path))

    if save_path.stat().st_size == 0:
        save_path.unlink()
        return jsonify({"error": "Uploaded file is empty."}), 400

    try:
        full_text, page_map = extract_file(save_path)
    except Exception as e:
        save_path.unlink(missing_ok=True)
        return jsonify({"error": f"Could not read file: {str(e)}"}), 422

    if not full_text.strip():
        save_path.unlink(missing_ok=True)
        return jsonify({"error": "File has no readable text (possibly scanned image PDF)."}), 422

    chunks = chunk_text(full_text, page_map, file.filename, chat_id)
    ids = [c["chunk_id"] for c in chunks]
    texts = [c["text"] for c in chunks]
    metadatas = [{"filename": c["filename"], "page": str(c["page"])} for c in chunks]

    collection = get_chat_collection(chat_id)
    try:
        existing = collection.get(where={"filename": file.filename})
        if existing["ids"]:
            collection.delete(ids=existing["ids"])
    except Exception:
        pass

    try:
        add_chunks_in_batches(collection, ids, texts, metadatas)
    except RuntimeError as e:
        # clean up so a failed upload doesn't leave an orphaned file on disk
        save_path.unlink(missing_ok=True)
        return jsonify({"error": str(e)}), 502

    record_file(chat_id, file.filename, stored_name, len(chunks), len(page_map))
    touch_chat(chat_id)

    return jsonify({
        "success": True,
        "filename": file.filename,
        "chunks": len(chunks),
        "pages": len(page_map),
    })


# ═══════════════════════════════════════════════════════
# ROUTES — FILES (list / delete / reindex) — scoped to one chat
# ═══════════════════════════════════════════════════════
@app.route("/api/chats/<chat_id>/files", methods=["GET"])
@login_required
@chat_access_required
def list_files(chat_id):
    files = list_chat_files(chat_id)
    return jsonify({"files": [{"name": f["filename"], "chunks": f["chunks"], "pages": f["pages"]} for f in files]})


@app.route("/api/chats/<chat_id>/files/<filename>", methods=["DELETE"])
@login_required
@chat_access_required
def delete_file(chat_id, filename):
    collection = get_chat_collection(chat_id)
    try:
        existing = collection.get(where={"filename": filename})
        if existing["ids"]:
            collection.delete(ids=existing["ids"])
    except Exception:
        pass

    record = get_chat_file(chat_id, filename)
    if record:
        (UPLOAD_DIR / record["stored_name"]).unlink(missing_ok=True)
    remove_chat_file(chat_id, filename)
    return jsonify({"success": True})


@app.route("/api/chats/<chat_id>/files/<filename>/reindex", methods=["POST"])
@login_required
@chat_access_required
def reindex_file(chat_id, filename):
    record = get_chat_file(chat_id, filename)
    if not record:
        return jsonify({"error": "File not found for this chat."}), 404

    save_path = UPLOAD_DIR / record["stored_name"]
    if not save_path.exists():
        return jsonify({"error": "File not found on disk. Please re-upload."}), 404

    try:
        full_text, page_map = extract_file(save_path)
        chunks = chunk_text(full_text, page_map, filename, chat_id)
        ids = [c["chunk_id"] for c in chunks]
        texts = [c["text"] for c in chunks]
        metadatas = [{"filename": c["filename"], "page": str(c["page"])} for c in chunks]

        collection = get_chat_collection(chat_id)
        try:
            existing = collection.get(where={"filename": filename})
            if existing["ids"]:
                collection.delete(ids=existing["ids"])
        except Exception:
            pass
        add_chunks_in_batches(collection, ids, texts, metadatas)

        record_file(chat_id, filename, record["stored_name"], len(chunks), len(page_map))
        return jsonify({"success": True, "chunks": len(chunks)})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ═══════════════════════════════════════════════════════
# ROUTES — CHAT HISTORY (scoped to one chat, persisted in SQLite)
# ═══════════════════════════════════════════════════════
@app.route("/api/chats/<chat_id>/history", methods=["GET"])
@login_required
@chat_access_required
def api_get_history(chat_id):
    history = get_history_for_display(chat_id, limit=MAX_HISTORY * 5)
    for m in history:
        img = m.pop("image_path", None)
        if img:
            m["image_url"] = f"/api/chats/{chat_id}/image/{img}"
    return jsonify({"history": history})


@app.route("/api/chats/<chat_id>/history", methods=["DELETE"])
@login_required
@chat_access_required
def api_clear_history(chat_id):
    clear_chat_history(chat_id)
    return jsonify({"success": True})


# ═══════════════════════════════════════════════════════
# ROUTES — CHAT (STREAMING) — scoped to one chat
# ═══════════════════════════════════════════════════════
@app.route("/api/chats/<chat_id>/chat", methods=["POST"])
@login_required
@chat_access_required
def chat_endpoint(chat_id):
    data = request.get_json(force=True)
    query = (data.get("message") or "").strip()
    if not query:
        return jsonify({"error": "Empty message"}), 400

    user = get_user(session["username"])
    user_id = user["id"]
    is_doc_request, doc_fmt = looks_like_document_request(query)

    if is_doc_request:
        return handle_chat_document_generation(chat_id, user_id, query, doc_fmt)
        
    files = list_chat_files(chat_id)
    chunks, sources = [], []
    if files:
        collection = get_chat_collection(chat_id)
        chunks = retrieve(collection, query)
        seen = set()
        for c in chunks:
            key = f"{c['filename']}|{c['page']}"
            if key not in seen:
                seen.add(key)
                sources.append({"filename": c["filename"], "page": c["page"]})

    memories = retrieve_memory(user_id, query)

    history = get_history(chat_id)
    messages = build_messages(history, query, chunks, memories)

    def generate():
        full_response = ""
        try:
            stream = ollama.chat(model=CHAT_MODEL, messages=messages, stream=True)
            for chunk in stream:
                token = chunk["message"]["content"]
                full_response += token
                yield f"data: {json.dumps({'type': 'token', 'content': token})}\n\n"

            if sources:
                yield f"data: {json.dumps({'type': 'sources', 'sources': sources})}\n\n"

            add_message(chat_id, "user", query)
            add_message(chat_id, "assistant", full_response)
            touch_chat(chat_id)

            yield f"data: {json.dumps({'type': 'done'})}\n\n"

            # Fire-and-forget: only actually runs (in a background thread) once
            # enough new turns have piled up since this chat's last extraction.
            maybe_trigger_memory_extraction(chat_id, user_id)
        except ollama.ResponseError as e:
            err = f"Model error: {e.error}. Is '{CHAT_MODEL}' pulled? Run: ollama pull {CHAT_MODEL}"
            yield f"data: {json.dumps({'type': 'error', 'content': err})}\n\n"
        except Exception as e:
            err = f"Unexpected error: {str(e)}"
            yield f"data: {json.dumps({'type': 'error', 'content': err})}\n\n"

    return Response(generate(), mimetype="text/event-stream",
                     headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})


VISION_SYSTEM_PROMPT = """You are Professor AI, looking at an image a faculty member has shared — this could be a scanned or handwritten student answer sheet, an exam/question paper, a lecture slide, a diagram, or a printed page from course material.

RULES:
1. Look at the image carefully before answering. If it contains handwriting, transcribe the relevant parts accurately; if something is genuinely illegible, say so rather than guessing silently and presenting the guess as fact.
2. If this looks like a student's answer sheet or submission and the professor is asking for evaluation, assess it against whatever criteria the professor gave (or general academic correctness if none given): note what's correct, what's missing or wrong, and give constructive, specific feedback. Do not assign a final numeric grade as if it were official unless the professor explicitly asks for a score — and if you do give one, frame it clearly as a suggested/draft assessment, not an authoritative record.
3. Never invent or assume a student's name, roll number, or identity beyond what's actually legible in the image.
4. If it's a diagram, slide, or printed material, describe and explain the underlying academic content, not just what's visually on the page.
5. MATH FORMATTING: Use LaTeX — \\( ... \\) inline, \\[ ... \\] for display equations.
6. Be professional, precise, and concise. No emotional or motivational commentary.
"""


@app.route("/api/chats/<chat_id>/chat-image", methods=["POST"])
@login_required
@chat_access_required
def chat_image_endpoint(chat_id):
    """Image analysis via a vision-capable Ollama model (e.g. Qwen2.5-VL) —
    for scanned/handwritten answer sheets, exam papers, slides, or diagrams.
    Optional: every other feature works fine if no vision model is pulled;
    this route will just return a clear error telling you to pull one."""
    image = request.files.get("image")
    query = (request.form.get("message") or "").strip()
    if not image or not image.filename:
        return jsonify({"error": "No image provided"}), 400

    ext = Path(image.filename).suffix.lower()
    if ext not in VISION_ALLOWED_EXT:
        return jsonify({"error": f"Unsupported image format '{ext}'. Use PNG, JPG, WEBP, GIF, or BMP."}), 400

    user = get_user(session["username"])
    user_id = user["id"]

    stored_name = f"chat{chat_id}_img_{int(time.time() * 1000)}{ext}"
    save_path = UPLOAD_DIR / stored_name
    image.save(str(save_path))

    size = save_path.stat().st_size if save_path.exists() else 0
    if size == 0:
        save_path.unlink(missing_ok=True)
        return jsonify({"error": "Image is empty."}), 400
    if size > VISION_MAX_BYTES:
        save_path.unlink(missing_ok=True)
        return jsonify({"error": f"Image is too large (max {VISION_MAX_BYTES // (1024 * 1024)} MB)."}), 400

    try:
        image_b64 = base64.b64encode(save_path.read_bytes()).decode("ascii")
    except Exception as e:
        save_path.unlink(missing_ok=True)
        return jsonify({"error": f"Could not read image: {e}"}), 422

    caption = query or "Analyze this image and summarize what it contains."
    memories = retrieve_memory(user_id, caption)
    history = get_history(chat_id)  # text-only prior turns — old images are never replayed

    system_content = VISION_SYSTEM_PROMPT
    if memories:
        system_content += _format_memory_block(memories)

    messages = [{"role": "system", "content": system_content}]
    messages.extend(history)
    messages.append({"role": "user", "content": caption, "images": [image_b64]})

    def generate():
        full_response = ""
        try:
            stream = ollama.chat(model=VISION_MODEL, messages=messages, stream=True,
                                  options={"num_predict": VISION_MAX_TOKENS})
            for chunk in stream:
                token = chunk["message"]["content"]
                full_response += token
                yield f"data: {json.dumps({'type': 'token', 'content': token})}\n\n"

            add_message(chat_id, "user", query or "[Image uploaded]", image_path=stored_name)
            add_message(chat_id, "assistant", full_response)
            touch_chat(chat_id)

            yield f"data: {json.dumps({'type': 'done'})}\n\n"

            maybe_trigger_memory_extraction(chat_id, user_id)
        except ollama.ResponseError as e:
            err = (f"Model error: {e.error}. Is '{VISION_MODEL}' pulled? "
                   f"Run: ollama pull {VISION_MODEL} — or skip image analysis if you don't need it.")
            yield f"data: {json.dumps({'type': 'error', 'content': err})}\n\n"
        except Exception as e:
            err = f"Unexpected error: {str(e)}"
            yield f"data: {json.dumps({'type': 'error', 'content': err})}\n\n"

    return Response(generate(), mimetype="text/event-stream",
                     headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})


@app.route("/api/chats/<chat_id>/image/<filename>", methods=["GET"])
@login_required
@chat_access_required
def serve_chat_image(chat_id, filename):
    """Re-serves a previously uploaded chat image. The chat<id>_img_ prefix
    check plus chat_access_required means a professor can't fetch another
    chat's — or another professor's — image by guessing a filename."""
    if not filename.startswith(f"chat{chat_id}_img_") or not (UPLOAD_DIR / filename).exists():
        return jsonify({"error": "Not found"}), 404
    return send_from_directory(UPLOAD_DIR.resolve(), filename)


# ═══════════════════════════════════════════════════════
# ROUTES — LONG-TERM MEMORY (scoped to the user, spans all their chats)
# ═══════════════════════════════════════════════════════
@app.route("/api/memory", methods=["GET"])
@login_required
def api_list_memory():
    user = get_user(session["username"])
    return jsonify({"memories": list_user_memory(user["id"]), "categories": list(MEMORY_CATEGORIES)})


@app.route("/api/memory", methods=["POST"])
@login_required
def api_add_memory():
    """Lets a professor directly teach Professor AI an academic preference —
    e.g. "I prefer scenario-based MCQs" — the same way you'd tell an
    assistant "remember that I prefer...".

    Unlike a general assistant, this does NOT accept free-form "remember
    anything" text: the spec restricts Professor AI's memory to academic
    preferences only, so the caller must pick one of MEMORY_CATEGORIES and
    store_memory_fact() will reject anything else outright (see its
    docstring). On top of that category restriction, a lightweight keyword
    guard blocks anything that looks like it references a specific
    student's data (name-linked marks, grades, attendance, rank, etc.) —
    that class of information must never enter long-term memory regardless
    of what category the caller claims it belongs to."""
    user = get_user(session["username"])
    data = request.get_json(force=True)
    text = (data.get("text") or "").strip()
    category = (data.get("category") or "").strip().lower()
    if not text:
        return jsonify({"error": "Memory text is required."}), 400
    if len(text) > 500:
        return jsonify({"error": "Keep a single memory under 500 characters."}), 400
    if category not in MEMORY_CATEGORIES:
        return jsonify({"error": f"category must be one of: {', '.join(MEMORY_CATEGORIES)}"}), 400
    if _STUDENT_DATA_GUARD_RE.search(text):
        return jsonify({"error": "This looks like it may reference student data (grades, "
                                  "attendance, scores, etc.), which Professor AI never stores "
                                  "as long-term memory. Rephrase as a general teaching "
                                  "preference if that's what you meant."}), 400

    # A fact the professor typed directly is as reliable as it gets — high
    # importance and full confidence, no guessing needed the way extraction
    # has to guess.
    mem_id = store_memory_fact(user["id"], text, category=category, source_chat_id=None,
                                dedup=True, importance=0.9, confidence=1.0)
    if mem_id is None:
        return jsonify({"error": "That's already stored (or too similar to an existing memory)."}), 409
    return jsonify({"success": True, "id": mem_id})


@app.route("/api/memory/<int:memory_id>", methods=["DELETE"])
@login_required
def api_delete_memory(memory_id):
    user = get_user(session["username"])
    ok = delete_memory_fact(user["id"], memory_id)
    if not ok:
        return jsonify({"error": "Memory not found."}), 404
    return jsonify({"success": True})


@app.route("/api/memory", methods=["DELETE"])
@login_required
def api_clear_memory():
    user = get_user(session["username"])
    clear_user_memory(user["id"])
    return jsonify({"success": True})


# ═══════════════════════════════════════════════════════
# ROUTES — DOCUMENT GENERATION ENGINE (Phase 1)
# ═══════════════════════════════════════════════════════
@app.route("/api/generate-document", methods=["POST"])
@login_required
def api_generate_document():
    data = request.get_json(force=True)
    prompt = (data.get("prompt") or "").strip()
    fmt = (data.get("format") or "pdf").strip().lower()

    if not prompt:
        return jsonify({"error": "A document prompt is required."}), 400
    if fmt not in DOC_ALLOWED_FORMATS:
        return jsonify({"error": f"format must be one of {list(DOC_ALLOWED_FORMATS)}"}), 400
    if fmt in ("docx", "both") and not DOCX_OK:
        return jsonify({"error": "python-docx is not installed on the server. Run: pip install python-docx"}), 500
    if fmt in ("pdf", "both") and not PDF_OK:
        return jsonify({"error": "reportlab is not installed on the server. Run: pip install reportlab"}), 500

    user = get_user(session["username"])
    user_id = user["id"]

    # Step 2: requirement extraction
    requirements = extract_document_requirements(prompt)

    # Step 3: AI content generation (structure only, not a file yet)
    try:
        structured = generate_document_content(
            requirements["title"], requirements["type"], requirements["length"], prompt
        )
    except ollama.ResponseError as e:
        return jsonify({"error": f"Model error: {e.error}. Is '{DOC_GEN_MODEL}' pulled?"}), 500
    except Exception as e:
        return jsonify({"error": f"Could not generate document content: {e}"}), 500

    # Step 4: render into the requested file format(s)
    doc_dir = get_user_doc_dir(user_id)
    base_name = make_safe_filename(structured["title"])

    files = []
    try:
        if fmt in ("pdf", "both"):
            pdf_filename = f"{base_name}.pdf"
            render_pdf(structured, doc_dir / pdf_filename)
            record_generated_document(user_id, structured["title"], requirements["type"], pdf_filename, "pdf")
            files.append({"filename": pdf_filename, "format": "pdf", "download_url": f"/download/{pdf_filename}"})
        if fmt in ("docx", "both"):
            docx_filename = f"{base_name}.docx"
            render_docx(structured, doc_dir / docx_filename)
            record_generated_document(user_id, structured["title"], requirements["type"], docx_filename, "docx")
            files.append({"filename": docx_filename, "format": "docx", "download_url": f"/download/{docx_filename}"})
    except Exception as e:
        return jsonify({"error": f"Document rendering failed: {e}"}), 500

    if fmt == "both":
        return jsonify({"success": True, "title": structured["title"], "files": files})

    # single-format request -> match the documented response shape exactly
    f = files[0]
    return jsonify({"success": True, "filename": f["filename"], "download_url": f["download_url"]})


@app.route("/api/documents", methods=["GET"])
@login_required
def api_list_documents():
    """Lists everything this user has generated, so the UI can show past
    documents instead of them only being reachable from the chat that made them."""
    user = get_user(session["username"])
    return jsonify({"documents": list_generated_documents(user["id"])})


@app.route("/download/<path:filename>", methods=["GET"])
@login_required
def download_document(filename):
    user = get_user(session["username"])
    doc_dir = get_user_doc_dir(user["id"])
    # Path(...).name strips any directory components the client might sneak
    # into the URL, so a request can never escape this user's own folder.
    safe_name = Path(filename).name
    file_path = doc_dir / safe_name
    if not file_path.exists() or not file_path.is_file():
        return jsonify({"error": "File not found."}), 404
    return send_from_directory(str(doc_dir), safe_name, as_attachment=True)


# ═══════════════════════════════════════════════════════
# SERVE FRONTEND
# ═══════════════════════════════════════════════════════
@app.route("/")
@login_required
def index():
    return render_template("index.html", name=session.get("name", ""))


# ═══════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════
if __name__ == "__main__":
    PORT = 5002  # different from the Student AI app's 5000, so both can run at once
    print("=" * 55)
    print(" Professor AI (per-chat isolation, faculty memory scope)")
    print(f" http://localhost:{PORT}")
    print("=" * 55)
    print(f" Chat model  : {CHAT_MODEL}")
    print(f" Embed model : {EMBED_MODEL}")
    print(f" Vision model: {VISION_MODEL} (optional — used only by Attach-image)")
    print(f" Chunk size  : {CHUNK_SIZE} chars | Overlap: {CHUNK_OVERLAP}")
    print(f" Top-K       : {TOP_K}")
    print(f" Memory      : extract every {MEMORY_EXTRACTION_THRESHOLD} msgs | top-{MEMORY_TOP_K} recall "
          f"| categories: {', '.join(MEMORY_CATEGORIES)}")
    print(f" RAG         : pool={TOP_K * RAG_CANDIDATE_MULTIPLIER} chunks | lexical weight={RAG_LEXICAL_WEIGHT}")
    print(f" Doc engine  : DOCX={'ok' if DOCX_OK else 'MISSING (pip install python-docx)'} "
          f"| PDF={'ok' if PDF_OK else 'MISSING (pip install reportlab)'} "
          f"| PPTX read={'ok' if PPTX_OK else 'MISSING (pip install python-pptx)'} | store: {GENERATED_DOCS_DIR}")
    print(f" DB          : {DB_PATH}")
    print(f" Chroma dir  : {CHROMA_DIR}")
    print(f" Uploads dir : {UPLOAD_DIR}")
    print("=" * 55)
    app.run(host="0.0.0.0", port=PORT, debug=False, threaded=True)